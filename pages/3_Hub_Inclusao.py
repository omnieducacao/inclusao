# ==============================================================================
# HUB DE INCLUS√ÉO - VERS√ÉO MODULAR E ORGANIZADA
# ==============================================================================

# --- IMPORTS ---
import streamlit as st
from datetime import date, datetime
from zoneinfo import ZoneInfo
import os
import re
import base64
import json
import requests
import pandas as pd
from PIL import Image
from io import BytesIO
import omni_utils as ou
from omni_utils import get_icon, icon_title, get_icon_emoji

# Importa√ß√µes OpenAI e ML
from openai import OpenAI

# Importa√ß√µes para documentos
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt, Inches, RGBColor
try:
    from docx.oxml.ns import qn
except ImportError:
    # Fallback se qn n√£o estiver dispon√≠vel
    qn = lambda x: x, RGBColor
from docx.oxml.ns import qn
from fpdf import FPDF
from pypdf import PdfReader

# Importa√ß√µes UI
from streamlit_cropper import st_cropper
import omni_utils as ou

# ==============================================================================
# CONFIGURA√á√ÉO INICIAL
# ==============================================================================

# ‚úÖ set_page_config UMA VEZ S√ì, SEMPRE no topo
st.set_page_config(
    page_title="Omnisfera | Hub de Recursos",
    page_icon="omni_icone.png",
    layout="wide",
    initial_sidebar_state="collapsed",
)

APP_VERSION = "v150.0 (SaaS Design)"

# ‚úÖ UI lockdown (n√£o quebra se faltar)
try:
    from ui_lockdown import hide_streamlit_chrome_if_needed, hide_default_sidebar_nav
    hide_streamlit_chrome_if_needed()
    hide_default_sidebar_nav()
except Exception:
    pass

# ‚úÖ Header + Navbar (depois do page_config)
ou.render_omnisfera_header()
ou.render_navbar(active_tab="Hub de Recursos")
if "hub" not in ou.get_enabled_modules():
    st.warning("Este m√≥dulo est√° desabilitado para sua escola.")
    if st.button("Voltar ao In√≠cio"):
        st.switch_page("pages/0_Home.py")
    st.stop()
ou.inject_compact_app_css()

# Adiciona classe no body para cores espec√≠ficas das abas
st.markdown("<script>document.body.classList.add('page-teal');</script>", unsafe_allow_html=True)

# ==============================================================================
# AJUSTE FINO DE LAYOUT (ANTES DO HERO - PADRONIZADO)
# ==============================================================================
def forcar_layout_hub():
    st.markdown("""
        <style>
            /* 1. Remove o cabe√ßalho padr√£o do Streamlit e a linha colorida */
            header[data-testid="stHeader"] {
                visibility: hidden !important;
                height: 0px !important;
            }

            /* 2. Puxa todo o conte√∫do para cima (O SEGREDO EST√Å AQUI) - ESPECIFICIDADE M√ÅXIMA */
            body .main .block-container,
            body .block-container,
            .main .block-container,
            .block-container {
                padding-top: 0px !important; /* SEM espa√ßo entre navbar e hero */
                padding-bottom: 1rem !important;
                margin-top: 0px !important;
            }
            
            /* 3. Remove qualquer espa√ßamento do Streamlit */
            [data-testid="stVerticalBlock"],
            div[data-testid="stVerticalBlock"] > div:first-child,
            .main .block-container > div:first-child,
            .main .block-container > *:first-child {
                padding-top: 0px !important;
                margin-top: 0px !important;
            }
            
            /* 4. Remove espa√ßamento do stMarkdown que renderiza o hero */
            .main .block-container > div:first-child .stMarkdown {
                margin-top: 0px !important;
                padding-top: 0px !important;
            }
            
            /* 5. Hero card colado no menu - margin negativo (ajustado para n√£o ficar muito colado) */
            .mod-card-wrapper {
                margin-top: -96px !important; /* Puxa o hero para cima, mas n√£o tanto quanto antes */
                position: relative;
                z-index: 1;
            }
            
            /* 6. Esconde o menu hamb√∫rguer e rodap√© */
            #MainMenu {visibility: hidden;}
            footer {visibility: hidden;}
        </style>
    """, unsafe_allow_html=True)

# CHAME ESTA FUN√á√ÉO ANTES DO HERO CARD (igual ao PEI)
forcar_layout_hub()

# Cores dos hero cards (paleta vibrante)
ou.inject_hero_card_colors()
# CSS padronizado: abas (p√≠lulas), bot√µes, selects, etc.
ou.inject_unified_ui_css()

# ==============================================================================
# HERO - HUB DE INCLUS√ÉO
# ==============================================================================
hora = datetime.now(ZoneInfo("America/Sao_Paulo")).hour
saudacao = "Bom dia" if 5 <= hora < 12 else "Boa tarde" if 12 <= hora < 18 else "Boa noite"
USUARIO_NOME = st.session_state.get("usuario_nome", "Visitante").split()[0]
WORKSPACE_NAME = st.session_state.get("workspace_name", "Workspace")

st.markdown(f"""
<div class="mod-card-wrapper">
    <div class="mod-card-rect">
        <div class="mod-bar c-teal"></div>
        <div class="mod-icon-area bg-teal-soft">
            <i class="ri-rocket-fill"></i>
        </div>
        <div class="mod-content">
            <div class="mod-title">Hub de Recursos & Intelig√™ncia Artificial</div>
            <div class="mod-desc">
                {saudacao}, <strong>{USUARIO_NOME}</strong>! Acesse recursos, modelos e ferramentas de IA 
                para criar adapta√ß√µes e estrat√©gias personalizadas no workspace <strong>{WORKSPACE_NAME}</strong>.
            </div>
        </div>
    </div>
</div>
""", unsafe_allow_html=True)

# CSS espec√≠fico do hero card do Hub
st.markdown("""
<style>
    /* CARD HERO - PADR√ÉO */
    .mod-card-wrapper { 
        display: flex; 
        flex-direction: column; 
        margin-bottom: 4px; 
        /* margin-top j√° aplicado no forcar_layout_hub() - n√£o duplicar aqui */
        border-radius: 16px; 
        overflow: hidden; 
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.02); 
        position: relative;
        z-index: 1;
    }
    .mod-card-rect { 
        background: white; 
        border-radius: 16px 16px 0 0; 
        padding: 0; 
        border: 1px solid #E2E8F0; 
        border-bottom: none; 
        display: flex; 
        flex-direction: row; 
        align-items: center; 
        height: 130px !important; 
        width: 100%; 
        position: relative; 
        overflow: hidden; 
        transition: all 0.25s cubic-bezier(0.4, 0, 0.2, 1); 
    }
    .mod-card-rect:hover { 
        transform: translateY(-4px); 
        box-shadow: 0 12px 24px rgba(0, 0, 0, 0.08); 
        border-color: #CBD5E1; 
    }
    .mod-bar { 
        width: 6px; 
        height: 100%; 
        flex-shrink: 0; 
    }
    .mod-icon-area { 
        width: 90px; 
        height: 100%; 
        display: flex; 
        align-items: center; 
        justify-content: center; 
        font-size: 1.8rem; 
        flex-shrink: 0; 
        background: #FAFAFA !important; 
        border-right: 1px solid #F1F5F9; 
        transition: all 0.3s ease; 
    }
    .mod-card-rect:hover .mod-icon-area { 
        background: white !important;
        transform: scale(1.05) !important;
    }
    .mod-content { 
        flex-grow: 1; 
        padding: 0 24px; 
        display: flex; 
        flex-direction: column; 
        justify-content: center; 
        min-width: 0;
        align-items: flex-start;
    }
    .mod-title { 
        font-weight: 800; 
        font-size: 1.1rem; 
        color: #1E293B; 
        margin-bottom: 6px; 
        letter-spacing: -0.3px; 
        transition: color 0.2s; 
    }
    .mod-desc { 
        font-size: 0.8rem; 
        color: #64748B; 
        line-height: 1.4; 
        display: -webkit-box; 
        -webkit-line-clamp: 2; 
        -webkit-box-orient: vertical; 
        overflow: hidden; 
    }
    
    /* CORES ESPEC√çFICAS TEAL - Garantir que o √≠cone tenha cor correta */
    .c-teal { background: #0D9488 !important; }
    .bg-teal-soft {
        background: #CCFBF1 !important;
        color: #0D9488 !important;
    }
    .mod-icon-area i { color: inherit !important; }
    .bg-teal-soft i,
    .mod-icon-area.bg-teal-soft i,
    .mod-icon-area.bg-teal-soft i.ri-rocket-fill {
        color: #0D9488 !important;
        font-size: 1.8rem !important;
    }
    .mod-card-rect:hover .mod-title {
        color: #0D9488; /* Specific hover color */
    }
</style>
""", unsafe_allow_html=True)

# Espa√ßamento ap√≥s hero card
st.markdown("<div style='height:8px'></div>", unsafe_allow_html=True)

# ==============================================================================
# VERIFICA√á√ÉO DE ACESSO
# ==============================================================================
def verificar_acesso():
    """Verifica se o usu√°rio est√° autenticado"""
    if "autenticado" not in st.session_state or not st.session_state["autenticado"]:
        ou.render_acesso_bloqueado("Fa√ßa login na P√°gina Inicial para acessar o Hub de Recursos.")
    try:
        from ui.permissions import can_access
        if not can_access("hub"):
            ou.render_acesso_bloqueado(
                "Voc√™ n√£o tem permiss√£o para acessar o Hub de Recursos.",
                "Entre em contato com o respons√°vel pela escola.",
            )
    except Exception:
        pass

verificar_acesso()

# ==============================================================================
# CONSTANTES E DADOS GLOBAIS
# ==============================================================================

# Dados da Taxonomia de Bloom
TAXONOMIA_BLOOM = {
    "1. Lembrar (Memorizar)": ["Citar", "Definir", "Identificar", "Listar", "Nomear", "Reconhecer", "Recordar", "Relacionar", "Repetir", "Sublinhar"],
    "2. Entender (Compreender)": ["Classificar", "Descrever", "Discutir", "Explicar", "Expressar", "Identificar", "Localizar", "Narrar", "Reafirmar", "Reportar", "Resumir", "Traduzir"],
    "3. Aplicar": ["Aplicar", "Demonstrar", "Dramatizar", "Empregar", "Esbo√ßar", "Ilustrar", "Interpretar", "Operar", "Praticar", "Programar", "Usar"],
    "4. Analisar": ["Analisar", "Calcular", "Categorizar", "Comparar", "Contrastar", "Criticar", "Diferenciar", "Discriminar", "Distinguir", "Examinar", "Experimentar", "Testar"],
    "5. Avaliar": ["Argumentar", "Avaliar", "Defender", "Escolher", "Estimar", "Julgar", "Prever", "Selecionar", "Suportar", "Validar", "Valorizar"],
    "6. Criar": ["Compor", "Construir", "Criar", "Desenhar", "Desenvolver", "Formular", "Investigar", "Planejar", "Produzir", "Propor"]
}

# Lista de componentes curriculares padr√£o
DISCIPLINAS_PADRAO = ["Matem√°tica", "Portugu√™s", "Ci√™ncias", "Hist√≥ria", "Geografia", "Artes", "Educa√ß√£o F√≠sica", "Ingl√™s", "Filosofia", "Sociologia"]

# Mapeamento: label do components (DB) ‚Üí nome usado no Hub/DISCIPLINAS_PADRAO
COMPONENT_TO_HUB = {
    "L√≠ngua Portuguesa": "Portugu√™s",
    "Arte": "Artes",
    "L√≠ngua Inglesa": "Ingl√™s",
}
# Inverso para BNCC (Hub ‚Üí DB/BNCC)
HUB_TO_COMPONENT = {v: k for k, v in COMPONENT_TO_HUB.items()}


def _get_hub_componentes_filtrados():
    """
    Retorna (opcoes, mostrar_selector, default) para o selectbox de componente no Hub.
    Se professor leciona 1 componente: s√≥ esse, sem selector.
    Se leciona >1: selector com apenas esses.
    Se coordena√ß√£o/master/todos: lista completa, com selector.
    """
    try:
        from ui.permissions import get_member_from_session
        from services.members_service import get_member_components
        member = get_member_from_session()
        if not member or member.get("link_type") != "turma":
            return DISCIPLINAS_PADRAO, True, None
        mid = member.get("id")
        comps = get_member_components(mid) if mid else []
        if not comps:
            return DISCIPLINAS_PADRAO, True, None
        # Mapeia labels DB ‚Üí Hub
        mapped = [COMPONENT_TO_HUB.get(c, c) for c in comps]
        opcoes = [x for x in DISCIPLINAS_PADRAO if x in mapped]
        if not opcoes:
            opcoes = mapped if mapped else DISCIPLINAS_PADRAO
        if len(opcoes) == 1:
            return opcoes, False, opcoes[0]
        return opcoes, True, opcoes[0] if opcoes else None
    except Exception:
        return DISCIPLINAS_PADRAO, True, None


def _filtrar_disciplinas_por_membro(disciplinas_bncc: list) -> tuple:
    """
    Para formul√°rios com BNCC (disciplinas vindas do df). Retorna (opcoes_filtradas, mostrar_selector, default).
    """
    try:
        from ui.permissions import get_member_from_session
        from services.members_service import get_member_components
        member = get_member_from_session()
        if not member or member.get("link_type") != "turma":
            return disciplinas_bncc, True, None
        mid = member.get("id")
        comps = get_member_components(mid) if mid else []
        if not comps:
            return disciplinas_bncc, True, None
        # BNCC usa mesmo nome que components (Arte, L√≠ngua Portuguesa, Matem√°tica)
        opcoes = [d for d in disciplinas_bncc if d in comps]
        if not opcoes:
            opcoes = comps if comps else disciplinas_bncc
        if len(opcoes) == 1:
            return opcoes, False, opcoes[0]
        return opcoes, True, opcoes[0] if opcoes else None
    except Exception:
        return disciplinas_bncc, True, None


# Campos de Experi√™ncia (Educa√ß√£o Infantil)
CAMPOS_EXPERIENCIA_EI = [
    "O eu, o outro e o n√≥s",
    "Corpo, gestos e movimentos",
    "Tra√ßos, sons, cores e formas",
    "Escuta, fala, pensamento e imagina√ß√£o",
    "Espa√ßos, tempos, quantidades, rela√ß√µes e transforma√ß√µes"
]

# Metodologias para Plano de Aula
METODOLOGIAS = [
    "Aula Expositiva Dialogada", 
    "Metodologia Ativa", 
    "Aprendizagem Baseada em Problemas", 
    "Ensino H√≠brido",
    "Sala de Aula Invertida", 
    "Rota√ß√£o por Esta√ß√µes"
]

# T√©cnicas Ativas
TECNICAS_ATIVAS = [
    "Gamifica√ß√£o", 
    "Sala de Aula Invertida", 
    "Aprendizagem Baseada em Projetos (PBL)", 
    "Rota√ß√£o por Esta√ß√µes", 
    "Peer Instruction",
    "Estudo de Caso", 
    "Aprendizagem Cooperativa"
]

# Recursos dispon√≠veis
RECURSOS_DISPONIVEIS = [
    "Quadro/Giz", 
    "Projetor/Datashow", 
    "Lousa Digital", 
    "Tablets/Celulares", 
    "Internet", 
    "Materiais Maker (Papel, Cola, etc)", 
    "Jogos de Tabuleiro", 
    "Laborat√≥rio", 
    "Material Dourado",
    "Recursos de CAA", 
    "V√≠deos Educativos"
]

# ==============================================================================
# FUN√á√ïES DE UTILIDADE GERAL
# ==============================================================================

def extrair_dados_docx(uploaded_file):
    """Extrai texto e imagens de um arquivo DOCX"""
    uploaded_file.seek(0)
    imagens = []
    texto = ""
    try:
        doc = Document(uploaded_file)
        texto = "\n".join([p.text for p in doc.paragraphs if p.text.strip() != ""])
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                img_data = rel.target_part.blob
                if len(img_data) > 1024:
                    imagens.append(img_data)
    except Exception as e:
        print(f"Erro ao extrair DOCX: {e}")
    return texto, imagens

def sanitizar_imagem(image_bytes):
    """Converte e otimiza imagem"""
    try:
        img = Image.open(BytesIO(image_bytes)).convert("RGB")
        out = BytesIO()
        img.save(out, format="JPEG", quality=90)
        return out.getvalue()
    except Exception as e:
        print(f"Erro ao sanitizar imagem: {e}")
        return None

def baixar_imagem_url(url):
    """Baixa imagem a partir de uma URL"""
    try:
        resp = requests.get(url, timeout=10)
        if resp.status_code == 200:
            return BytesIO(resp.content)
    except Exception as e:
        print(f"Erro ao baixar imagem: {e}")
    return None

def buscar_imagem_unsplash(query, access_key):
    """Busca imagem no Unsplash"""
    if not access_key:
        return None
    url = f"https://api.unsplash.com/search/photos?query={query}&per_page=1&client_id={access_key}&lang=pt"
    try:
        resp = requests.get(url, timeout=5)
        data = resp.json()
        if data.get('results'):
            return data['results'][0]['urls']['regular']
    except Exception as e:
        print(f"Erro ao buscar no Unsplash: {e}")
    return None

def garantir_tag_imagem(texto):
    """Garante que o texto tenha pelo menos uma tag de imagem"""
    if "[[IMG" not in texto.upper() and "[[GEN_IMG" not in texto.upper():
        match = re.search(r'(\n|\. )', texto)
        if match:
            pos = match.end()
            return texto[:pos] + "\n\n[[IMG_1]]\n\n" + texto[pos:]
        return texto + "\n\n[[IMG_1]]"
    return texto

def construir_docx_final(texto_ia, aluno, materia, mapa_imgs, tipo_atv, sem_cabecalho=False, checklist_adaptacao=None):
    """Constr√≥i documento DOCX final com formata√ß√£o melhorada baseada no checklist"""
    doc = Document()
    
    # Determinar formata√ß√£o baseada no checklist
    usar_caixa_alta = False
    usar_opendyslexic = False
    espacamento_linhas = 1.5  # padr√£o
    
    if checklist_adaptacao:
        # Se precisa de adapta√ß√£o visual, usar formata√ß√£o especial
        if checklist_adaptacao.get("paragrafos_curtos") or not checklist_adaptacao.get("compreende_instrucoes_complexas"):
            usar_caixa_alta = True
            usar_opendyslexic = True
            espacamento_linhas = 1.8
    
    # Estilos personalizados
    style = doc.styles['Normal']
    if usar_opendyslexic:
        # Tentar usar OpenDyslexic, fallback para Arial se n√£o dispon√≠vel
        try:
            style.font.name = 'OpenDyslexic'
            # Tentar registrar fonte alternativa (pode n√£o funcionar se fonte n√£o estiver instalada)
            try:
                rFonts = style.element.rPr.rFonts
                if rFonts is not None:
                    rFonts.set(qn('w:ascii'), 'OpenDyslexic')
                    rFonts.set(qn('w:hAnsi'), 'OpenDyslexic')
            except:
                pass  # Se n√£o conseguir registrar, continua com o nome
        except:
            style.font.name = 'Arial'
    else:
        style.font.name = 'Arial'
    
    style.font.size = Pt(14)  # Aumentado para melhor legibilidade
    style.paragraph_format.space_after = Pt(8)
    style.paragraph_format.line_spacing = espacamento_linhas
    
    if not sem_cabecalho:
        # T√≠tulo principal
        titulo = doc.add_heading(f'{tipo_atv.upper()} ADAPTADA - {materia.upper()}', 0)
        titulo.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if titulo.runs:
            titulo.runs[0].font.size = Pt(16)
            titulo.runs[0].bold = True
        
        # Informa√ß√µes do estudante
        p_estudante = doc.add_paragraph(f"Estudante: {aluno['nome']}")
        p_estudante.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if p_estudante.runs:
            p_estudante.runs[0].font.size = Pt(11)
        
        # Linha separadora
        linha_sep = doc.add_paragraph("_"*50)
        linha_sep.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # T√≠tulo de se√ß√£o
        secao = doc.add_heading('Atividades', level=2)
        if secao.runs:
            secao.runs[0].font.size = Pt(14)

    linhas = texto_ia.split('\n')
    for linha in linhas:
        linha_limpa = linha.strip()
        if not linha_limpa:
            # Linha vazia - adiciona espa√ßo
            doc.add_paragraph()
            continue
            
        tag_match = re.search(r'\[\[(IMG|GEN_IMG).*?(\d+)\]\]', linha, re.IGNORECASE)
        if tag_match:
            partes = re.split(r'(\[\[(?:IMG|GEN_IMG).*?\d+\]\])', linha, flags=re.IGNORECASE)
            for parte in partes:
                sub_match = re.search(r'(\d+)', parte)
                if ("IMG" in parte.upper() or "GEN_IMG" in parte.upper()) and sub_match:
                    num = int(sub_match.group(1))
                    img_bytes = mapa_imgs.get(num)
                    if not img_bytes and len(mapa_imgs) == 1:
                        img_bytes = list(mapa_imgs.values())[0]
                    if img_bytes:
                        try:
                            p = doc.add_paragraph()
                            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            r = p.add_run()
                            r.add_picture(BytesIO(img_bytes), width=Inches(4.5))
                            # Espa√ßo ap√≥s imagem
                            doc.add_paragraph()
                        except Exception as e:
                            print(f"Erro ao adicionar imagem: {e}")
                elif parte.strip():
                    # Formatar texto com t√≠tulos e listas (passando par√¢metros de formata√ß√£o)
                    _adicionar_paragrafo_formatado(doc, parte.strip(), usar_caixa_alta, usar_opendyslexic, espacamento_linhas)
        else:
            # Formatar texto com t√≠tulos e listas (passando par√¢metros de formata√ß√£o)
            _adicionar_paragrafo_formatado(doc, linha_limpa, usar_caixa_alta, usar_opendyslexic, espacamento_linhas)
            
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def _adicionar_paragrafo_formatado(doc, texto, usar_caixa_alta=False, usar_opendyslexic=False, espacamento=1.5):
    """Adiciona par√°grafo formatado com detec√ß√£o de t√≠tulos e listas"""
    try:
        # Aplicar caixa alta se necess√°rio
        texto_formatado = texto.upper() if usar_caixa_alta else texto
        
        # Detectar t√≠tulos (come√ßam com # ou s√£o n√∫meros seguidos de ponto)
        if re.match(r'^#{1,3}\s+', texto_formatado):
            nivel = len(re.match(r'^(#+)', texto_formatado).group(1))
            texto_limpo = re.sub(r'^#+\s+', '', texto_formatado)
            heading = doc.add_heading(texto_limpo, level=min(nivel, 3))
            if heading.runs:
                heading.runs[0].font.size = Pt(16 - nivel)
                if usar_opendyslexic:
                    heading.runs[0].font.name = 'OpenDyslexic'
        elif re.match(r'^\d+[\.\)]\s+', texto_formatado):
            # Lista numerada
            p = doc.add_paragraph(texto_formatado, style='List Number')
            if p.runs:
                p.runs[0].font.size = Pt(14)
                p.runs[0].font.name = 'OpenDyslexic' if usar_opendyslexic else 'Arial'
                p.paragraph_format.line_spacing = espacamento
        elif re.match(r'^[-‚Ä¢*]\s+', texto_formatado):
            # Lista com marcadores
            texto_limpo = re.sub(r'^[-‚Ä¢*]\s+', '', texto_formatado)
            p = doc.add_paragraph(texto_limpo, style='List Bullet')
            if p.runs:
                p.runs[0].font.size = Pt(14)
                p.runs[0].font.name = 'OpenDyslexic' if usar_opendyslexic else 'Arial'
                p.paragraph_format.line_spacing = espacamento
        elif texto_formatado.isupper() and len(texto_formatado) < 100:
            # T√≠tulo em mai√∫sculas
            p = doc.add_paragraph(texto_formatado)
            if p.runs:
                p.runs[0].font.size = Pt(15)
                p.runs[0].bold = True
                p.runs[0].font.name = 'OpenDyslexic' if usar_opendyslexic else 'Arial'
                p.paragraph_format.line_spacing = espacamento
        else:
            # Texto normal
            p = doc.add_paragraph(texto_formatado)
            if p.runs:
                p.runs[0].font.size = Pt(14)
                p.runs[0].font.name = 'OpenDyslexic' if usar_opendyslexic else 'Arial'
                p.paragraph_format.line_spacing = espacamento
    except Exception as e:
        # Fallback: adiciona como par√°grafo simples
        doc.add_paragraph(texto_formatado if usar_caixa_alta else texto)

def criar_docx_simples(texto, titulo="Documento"):
    """Cria um DOCX simples a partir de texto"""
    doc = Document()
    doc.add_heading(titulo, 0)
    for para in texto.split('\n'):
        if para.strip():
            doc.add_paragraph(para.strip())
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def criar_pdf_generico(texto):
    """Cria um PDF simples a partir de texto"""
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    texto_safe = texto.encode('latin-1', 'replace').decode('latin-1')
    pdf.multi_cell(0, 10, texto_safe)
    return pdf.output(dest='S').encode('latin-1')


def gerar_ppt_do_plano_kimi(texto_plano: str, titulo_plano: str, aluno: dict = None, kimi_key: str = None) -> tuple[bytes | None, str | None]:
    """
    Gera PowerPoint rico visualmente a partir do plano de aula, usando Omnisfera Green (Kimi).
    Est√©tica inspirada no hiperfoco do aluno. Para o professor usar em aula.
    Retorna (bytes_pptx, None) em sucesso ou (None, mensagem_erro).
    """
    if not kimi_key:
        kimi_key = (
            st.session_state.get("_kimi_key_cached")
            or os.environ.get("OPENROUTER_API_KEY")
            or os.environ.get("KIMI_API_KEY")
            or ou.get_setting("OPENROUTER_API_KEY", "")
            or ou.get_setting("KIMI_API_KEY", "")
            or st.session_state.get("OPENROUTER_API_KEY", "")
            or st.session_state.get("KIMI_API_KEY", "")
        )
        try:
            kimi_key = kimi_key or (st.secrets.get("OPENROUTER_API_KEY", "") or st.secrets.get("KIMI_API_KEY", "") or "")
        except Exception:
            pass
        kimi_key = (kimi_key or "").strip() or None
    if not kimi_key:
        return None, (
            "Configure OPENROUTER_API_KEY ou KIMI_API_KEY em .streamlit/secrets.toml ou vari√°veis de ambiente. "
            "Chave OpenRouter (sk-or-...) em openrouter.ai/keys. "
            "Em deploy (Streamlit Cloud/Render): adicione nas Secrets/Env Vars do servi√ßo."
        )

    def _clean(s):
        return (s or "").strip().replace("\n", "").replace("\r", "").strip()

    kimi_key = _clean(kimi_key)

    # Usa OPENROUTER_BASE_URL, KIMI_MODEL se configurados; sen√£o detecta por prefixo da chave
    base_url = (
        os.environ.get("OPENROUTER_BASE_URL")
        or ou.get_setting("OPENROUTER_BASE_URL", "")
        or ("https://openrouter.ai/api/v1" if kimi_key.startswith("sk-or-") else "https://api.moonshot.ai/v1")
    )
    base_url = _clean(base_url) or "https://openrouter.ai/api/v1"
    model = (
        os.environ.get("KIMI_MODEL")
        or ou.get_setting("KIMI_MODEL", "")
        or ("moonshotai/kimi-k2.5" if kimi_key.startswith("sk-or-") else "kimi-k2-turbo-preview")
    )
    model = _clean(model) or "moonshotai/kimi-k2.5"

    hiperfoco = (aluno or {}).get("hiperfoco", "") or "aprendizado e descobertas"
    nome_aluno = (aluno or {}).get("nome", "") or "estudante"

    prompt = f"""Voc√™ √© um designer de apresenta√ß√µes profissionais. Crie uma APRESENTA√á√ÉO IMPACTANTE para o professor usar em sala ‚Äî n√≠vel executivo/corporativo, n√£o amador.

CONTEXTO:
- Plano de aula: {titulo_plano[:80]}
- Hiperfoco do aluno (est√©tica/theme): {hiperfoco[:200]}
- Nome do aluno: {nome_aluno}

OBJETIVO: PPT profissional, engajante, com hierarquia visual clara. Use o plano de aula para extrair os CONCEITOS-CHAVE, OBJETIVOS, ATIVIDADES e DESTAQUES. Cada slide deve ter prop√≥sito ‚Äî n√£o preencha com texto gen√©rico.

Retorne APENAS um JSON v√°lido, sem markdown:
{{
  "theme": "tema visual baseado no hiperfoco (ex: explora√ß√£o espacial, floresta, dinossauros, jogos)",
  "palette": {{
    "primary": "hex sem # (cor principal)",
    "accent": "hex sem # (cor de destaque)",
    "background_light": "hex (fundo claro)",
    "text_dark": "hex (texto principal)",
    "text_muted": "hex (texto secund√°rio)"
  }},
  "slides": [
    {{"slide_type": "cover", "title": "T√≠tulo da aula", "subtitle": "Frase de impacto com foco no hiperfoco"}},
    {{"slide_type": "section", "title": "Objetivos", "statement": "Frase curta e poderosa"}},
    {{"slide_type": "content", "title": "Conceito 1", "bullets": ["Bullet curto", "Outro bullet"], "callout": "Frase para destacar em caixa"}},
    {{"slide_type": "quote", "quote_text": "Frase inspiradora do plano", "quote_author": "contexto"}},
    {{"slide_type": "key_takeaway", "title": "O que levar", "bullets": ["Ponto 1", "Ponto 2"], "highlight": [0]}}
  ]
}}

TIPOS DE SLIDE:
- cover: t√≠tulo grande + subtitle
- section: divisor de se√ß√£o, statement grande (1 frase)
- content: t√≠tulo + bullets + callout opcional (frase em destaque)
- quote: cita√ß√£o centralizada, impactante
- key_takeaway: resumo com bullets, highlight = √≠ndices dos bullets mais importantes (0-based)

REGRAS:
- 8 a 14 slides
- Use o PLANO DE AULA para extrair conte√∫do real ‚Äî objetivos, metodologia, avalia√ß√£o
- T√≠tulos curtos e impactantes (m√°x 8 palavras)
- Bullets concisos (m√°x 50 chars cada)
- Cores que combinem com o tema/hiperfoco
- callout: frase curta para caixa de destaque (opcional)
- highlight: array de √≠ndices (ex: [0,2]) para bullets a enfatizar

--- PLANO DE AULA ---
{texto_plano[:5500]}
--- FIM ---"""

    try:
        client = OpenAI(
            api_key=kimi_key,
            base_url=base_url,
        )
        resp = client.chat.completions.create(
            model=model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
        )
        raw = (resp.choices[0].message.content or "").strip()
        if not raw:
            return None, "Resposta vazia do Kimi."

        # Extrair JSON (pode vir em ```json ... ```)
        if "```" in raw:
            for sep in ("```json", "```"):
                if sep in raw:
                    idx = raw.find(sep) + len(sep)
                    raw = raw[idx:].split("```")[0].strip()
                    break
        data = json.loads(raw)
        slides_data = data.get("slides") or []
        if not slides_data:
            return None, "Kimi n√£o retornou slides estruturados."

        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor as PptxRgb
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN

        def _hex_to_rgb(hex_str):
            if not hex_str or not isinstance(hex_str, str):
                return (15, 118, 110)  # teal default
            h = str(hex_str).strip().lstrip("#")[:6]
            if len(h) != 6:
                return (15, 118, 110)
            try:
                return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))
            except Exception:
                return (15, 118, 110)

        palette = data.get("palette") or {}
        rgb_primary = _hex_to_rgb(palette.get("primary") or palette.get("primary_rgb") or "0F766E")
        rgb_accent = _hex_to_rgb(palette.get("accent") or palette.get("accent_rgb") or "F59E0B")
        rgb_bg = _hex_to_rgb(palette.get("background_light") or "F0FDFA")
        rgb_text = _hex_to_rgb(palette.get("text_dark") or "0F172A")
        rgb_muted = _hex_to_rgb(palette.get("text_muted") or "64748B")
        clr_primary = PptxRgb(*rgb_primary)
        clr_accent = PptxRgb(*rgb_accent)
        clr_bg = PptxRgb(*rgb_bg)
        clr_text = PptxRgb(*rgb_text)
        clr_text_light = PptxRgb(*rgb_muted)
        margin = Inches(0.6)
        content_w = Inches(8.8)
        content_left = margin

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        for i, s in enumerate(slides_data):
            slide_type = str(s.get("slide_type") or "content").lower()
            title = str(s.get("title") or f"Slide {i+1}")[:200]
            subtitle = str(s.get("subtitle") or "")[:150]
            statement = str(s.get("statement") or "")[:200]
            quote_text = str(s.get("quote_text") or "")[:300]
            quote_author = str(s.get("quote_author") or "")[:80]
            callout = str(s.get("callout") or "")[:120]
            bullets = s.get("bullets") or []
            if isinstance(bullets, str):
                bullets = [bullets] if bullets.strip() else []
            bullets = [str(b)[:100] for b in bullets[:6]]
            highlight = s.get("highlight") or []
            if not isinstance(highlight, (list, tuple)):
                highlight = []
            highlight_set = set(int(x) for x in highlight if isinstance(x, (int, float)))

            layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(layout)
            bg = slide.background
            bg.fill.solid()
            bg.fill.fore_color.rgb = clr_bg

            if slide_type == "cover":
                bar_h = Inches(0.12)
                bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, margin, Inches(2.2), content_w, bar_h)
                bar.fill.solid()
                bar.fill.fore_color.rgb = clr_primary
                bar.line.fill.background()
                tx = slide.shapes.add_textbox(margin, Inches(0.8), content_w, Inches(1.8))
                tf = tx.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(44)
                p.font.bold = True
                for r in p.runs:
                    r.font.color.rgb = clr_primary
                if subtitle:
                    p2 = tf.add_paragraph()
                    p2.text = subtitle
                    p2.font.size = Pt(22)
                    p2.font.italic = True
                    for r in p2.runs:
                        r.font.color.rgb = clr_text_light
                    p2.space_before = Pt(12)
            elif slide_type == "section":
                bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, margin, Inches(0.4), Inches(0.08), Inches(6.2))
                bar.fill.solid()
                bar.fill.fore_color.rgb = clr_primary
                bar.line.fill.background()
                tx = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(7.6), Inches(2.5))
                tf = tx.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = statement or title
                p.font.size = Pt(36)
                p.font.bold = True
                for r in p.runs:
                    r.font.color.rgb = clr_primary
                if title and not statement:
                    p2 = tf.add_paragraph()
                    p2.text = title
                    p2.font.size = Pt(28)
                    for r in p2.runs:
                        r.font.color.rgb = clr_text
            elif slide_type == "quote":
                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(2.8))
                box.fill.solid()
                box.fill.fore_color.rgb = PptxRgb(248, 250, 252)
                box.line.color.rgb = clr_primary
                box.line.width = Pt(2)
                tx = box.text_frame
                tx.word_wrap = True
                p = tx.paragraphs[0]
                p.text = f'"{quote_text}"'
                p.font.size = Pt(26)
                p.font.italic = True
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.color.rgb = clr_text
                if quote_author:
                    p2 = tx.add_paragraph()
                    p2.text = f"‚Äî {quote_author}"
                    p2.font.size = Pt(16)
                    p2.alignment = PP_ALIGN.CENTER
                    for r in p2.runs:
                        r.font.color.rgb = clr_text_light
            else:
                bar_h = Inches(0.06)
                bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, margin, Inches(0.35), content_w, bar_h)
                bar.fill.solid()
                bar.fill.fore_color.rgb = clr_primary
                bar.line.fill.background()
                top = Inches(0.6)
                tx = slide.shapes.add_textbox(content_left, top, content_w, Inches(1))
                tf = tx.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(28)
                p.font.bold = True
                for r in p.runs:
                    r.font.color.rgb = clr_primary
                top = Inches(1.4)
                if bullets:
                    tx_body = slide.shapes.add_textbox(content_left, top, content_w, Inches(3.8))
                    tf_body = tx_body.text_frame
                    tf_body.word_wrap = True
                    for j, b in enumerate(bullets):
                        if j == 0:
                            p_b = tf_body.paragraphs[0]
                        else:
                            p_b = tf_body.add_paragraph()
                        p_b.text = f"‚Ä¢ {b}"
                        p_b.font.size = Pt(20 if j in highlight_set else 18)
                        p_b.font.bold = (j in highlight_set)
                        for r in p_b.runs:
                            r.font.color.rgb = clr_primary if j in highlight_set else clr_text
                        p_b.space_after = Pt(12)
                callout_top = Inches(5.2) if bullets else Inches(1.5)
                if callout:
                    call_h = Inches(0.9)
                    call_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, content_left, callout_top, content_w, call_h)
                    call_box.fill.solid()
                    call_box.fill.fore_color.rgb = clr_primary
                    call_box.line.fill.background()
                    tx_call = call_box.text_frame
                    tx_call.word_wrap = True
                    p_call = tx_call.paragraphs[0]
                    p_call.text = callout
                    p_call.font.size = Pt(18)
                    p_call.font.bold = True
                    for r in p_call.runs:
                        r.font.color.rgb = PptxRgb(255, 255, 255)
                visual_cue = str(s.get("visual_cue") or "")[:80]
                if visual_cue:
                    tx_cue = slide.shapes.add_textbox(content_left, Inches(6.5), content_w, Inches(0.5))
                    tf_cue = tx_cue.text_frame
                    p_cue = tf_cue.paragraphs[0]
                    p_cue.text = f"üñº {visual_cue}"
                    p_cue.font.size = Pt(12)
                    p_cue.font.italic = True
                    for r in p_cue.runs:
                        r.font.color.rgb = clr_text_light

        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read(), None
    except json.JSONDecodeError as e:
        return None, f"Resposta do Kimi n√£o √© JSON v√°lido: {str(e)[:100]}"
    except Exception as e:
        err = str(e)
        if "401" in err or "Invalid Authentication" in err or "invalid_authentication" in err:
            return None, (
                "Chave da API inv√°lida (401). Verifique: "
                "1) KIMI_API_KEY ou MOONSHOT_API_KEY em secrets.toml ou vari√°vel de ambiente; "
                "2) Chave correta de platform.moonshot.ai (n√£o use chave da OpenAI); "
                "3) Conta com saldo e chave ativa."
            )
        return None, err[:200]


# ==============================================================================
# FUN√á√ïES DE CONEX√ÉO COM SUPABASE (CORRIGIDO PARA SEPARAR DIAGN√ìSTICO DE HIPERFOCO)
# ==============================================================================
# Fun√ß√µes _sb_url(), _sb_key(), _headers() removidas - usar ou._sb_url(), ou._sb_key(), ou._headers() do omni_utils

@st.cache_data(ttl=10, show_spinner=False)
def list_students_rest(workspace_id: str = ""):
    """
    Busca estudantes do Supabase.
    workspace_id como argumento evita cache incorreto (cache por workspace).
    """
    if not (workspace_id and str(workspace_id).strip()):
        return []

    try:
        base = (
            f"{ou._sb_url()}/rest/v1/students"
            f"?select=id,name,grade,class_group,diagnosis,created_at,pei_data,paee_ciclos,planejamento_ativo"
            f"&workspace_id=eq.{workspace_id}"
            f"&order=created_at.desc"
        )
        r = requests.get(base, headers=ou._headers(), timeout=20)
        return r.json() if r.status_code == 200 else []
    except Exception as e:
        print(f"Erro ao carregar estudantes: {str(e)}")
        return []

def carregar_estudantes_supabase():
    """Carrega e processa estudantes, separando Diagn√≥stico de Hiperfoco. Filtra por membro se gest√£o ativa."""
    workspace_id = (st.session_state.get("workspace_id") or "").strip()
    if st.session_state.get("students_cache_invalid"):
        list_students_rest.clear()
        st.session_state.pop("students_cache_invalid", None)
    dados = list_students_rest(workspace_id)
    try:
        from ui.permissions import apply_member_filter
        dados = apply_member_filter(dados)
    except Exception:
        pass
    estudantes = []

    for item in dados:
        pei_completo = item.get("pei_data") or {}
        contexto_ia = ""

        # Tenta pegar texto de contexto da IA dentro do JSON do PEI
        if isinstance(pei_completo, dict):
            contexto_ia = pei_completo.get("ia_sugestao", "") or ""
        
        # 1. CAPTURA DO DIAGN√ìSTICO 
        # Primeiro tenta dentro do pei_data (onde o PEI salva como "diagnostico")
        diagnostico_real = ""
        if isinstance(pei_completo, dict):
            diagnostico_real = pei_completo.get("diagnostico") or ""
        
        # Se n√£o encontrou no pei_data, tenta da coluna do banco
        if not diagnostico_real:
            diagnostico_real = item.get("diagnosis") or ""
        
        # Se ainda n√£o encontrou, mostra mensagem
        if not diagnostico_real:
            diagnostico_real = "N√£o informado no cadastro"

        # 2. CAPTURA DO HIPERFOCO (Vem APENAS de dentro do JSON pei_data)
        # IMPORTANTE: N√ÉO usar diagn√≥stico como fallback para hiperfoco
        hiperfoco_real = ""
        if isinstance(pei_completo, dict):
            # Buscar hiperfoco - garantir que N√ÉO seja o diagn√≥stico
            hiperfoco_temp = (
                pei_completo.get("hiperfoco") or 
                pei_completo.get("interesses") or 
                pei_completo.get("habilidades_interesses") or 
                ""
            )
            # Se o hiperfoco encontrado for igual ao diagn√≥stico, descartar
            if hiperfoco_temp and hiperfoco_temp != diagnostico_real:
                hiperfoco_real = hiperfoco_temp
        
        # Se n√£o achou no JSON, coloca um padr√£o (NUNCA usa o diagn√≥stico)
        if not hiperfoco_real:
            hiperfoco_real = "Interesses gerais (A descobrir)"

        # Montagem do resumo de fallback se n√£o houver contexto da IA
        if not contexto_ia:
            serie = item.get("grade", "")
            contexto_ia = f"Estudante: {item.get('name')}. S√©rie: {serie}. Diagn√≥stico: {diagnostico_real}."

        estudante = {
            "nome": item.get("name", ""),
            "serie": item.get("grade", ""),
            
            # AQUI EST√Å A CORRE√á√ÉO: Chaves separadas e limpas
            "diagnosis": diagnostico_real,  
            "hiperfoco": hiperfoco_real,
            
            "ia_sugestao": contexto_ia,
            "id": item.get("id", ""),
            "pei_data": pei_completo,
        }
        
        if estudante["nome"]:
            estudantes.append(estudante)

    return estudantes

# ==============================================================================
# FUN√á√ïES DE IA (OPENAI)
# ==============================================================================

def gerar_imagem_inteligente(api_key, prompt, unsplash_key=None, feedback_anterior="", prioridade="IA", gemini_key=None):
    """
    Gera imagem: prioridade Gemini (se GEMINI_API_KEY), depois DALL-E, depois Unsplash.
    Retorna bytes (PNG) ou URL (str); st.image() aceita ambos.
    """
    # 1. TENTATIVA GEMINI (se chave configurada)
    key_gemini = gemini_key or ou.get_gemini_api_key()
    if key_gemini and prioridade == "IA":
        img_bytes, err = ou.gerar_imagem_ilustracao_gemini(prompt, feedback_anterior=feedback_anterior, api_key=key_gemini)
        if img_bytes:
            return img_bytes

    # 2. TENTATIVA BANCO DE IMAGENS (Se solicitado e configurado)
    if prioridade == "BANCO" and unsplash_key:
        termo = prompt.split('.')[0] if '.' in prompt else prompt
        url_banco = buscar_imagem_unsplash(termo, unsplash_key)
        if url_banco:
            return url_banco

    # 3. TENTATIVA OPENAI (DALL-E 3) ‚Äî s√≥ se houver chave v√°lida
    if not api_key or not str(api_key).strip():
        if unsplash_key and prioridade == "IA":
            termo = prompt.split('.')[0] if '.' in prompt else prompt
            return buscar_imagem_unsplash(termo, unsplash_key)
        return None
    try:
        client = OpenAI(api_key=api_key)
        prompt_final = f"{prompt}. Adjustment requested: {feedback_anterior}" if feedback_anterior else prompt
        didactic_prompt = f"Educational textbook illustration, clean flat vector style, white background. CRITICAL RULE: STRICTLY NO TEXT, NO TYPOGRAPHY, NO ALPHABET, NO NUMBERS, NO LABELS inside the image. Just the visual representation of: {prompt_final}"
        resp = client.images.generate(model="dall-e-3", prompt=didactic_prompt, size="1024x1024", quality="standard", n=1)
        return resp.data[0].url
    except Exception as e:
        if prioridade == "IA" and unsplash_key:
            termo = prompt.split('.')[0] if '.' in prompt else prompt
            return buscar_imagem_unsplash(termo, unsplash_key)
        print(f"Erro ao gerar imagem: {e}")
        return None

def gerar_pictograma_caa(api_key, conceito, feedback_anterior="", gemini_key=None):
    """
    Gera s√≠mbolo CAA: prioridade Gemini (se GEMINI_API_KEY), depois DALL-E.
    Retorna bytes (PNG) ou URL (str); st.image() aceita ambos.
    """
    key_gemini = gemini_key or ou.get_gemini_api_key()
    if key_gemini:
        img_bytes, err = ou.gerar_imagem_pictograma_caa_gemini(conceito, feedback_anterior=feedback_anterior, api_key=key_gemini)
        if img_bytes:
            return img_bytes
    if not api_key or not str(api_key).strip():
        return None
    try:
        client = OpenAI(api_key=api_key)
        ajuste = f" CORRE√á√ÉO PEDIDA: {feedback_anterior}" if feedback_anterior else ""
        prompt_caa = f"""
    Create a COMMUNICATION SYMBOL (AAC/PECS) for the concept: '{conceito}'. {ajuste}
    STYLE GUIDE:
    - Flat vector icon (ARASAAC/Noun Project style).
    - Solid WHITE background.
    - Thick BLACK outlines.
    - High contrast primary colors.
    - No background details, no shadows.
    - CRITICAL MANDATORY RULE: MUTE IMAGE. NO TEXT. NO WORDS. NO LETTERS. NO NUMBERS. 
    - The image must be a purely visual symbol.
    """
        resp = client.images.generate(model="dall-e-3", prompt=prompt_caa, size="1024x1024", quality="standard", n=1)
        return resp.data[0].url
    except Exception as e:
        print(f"Erro ao gerar pictograma: {e}")
        return None

def adaptar_conteudo_docx(api_key, aluno, texto, materia, tema, tipo_atv, remover_resp, questoes_mapeadas, modo_profundo=False, checklist_adaptacao=None):
    """Adapta conte√∫do de um DOCX para o estudante"""
    client = OpenAI(api_key=api_key)
    lista_q = ", ".join([str(n) for n in questoes_mapeadas]) if questoes_mapeadas else ""
    style = "Seja did√°tico e use uma Cadeia de Pensamento para adaptar." if modo_profundo else "Seja objetivo."
    
    # Montar instru√ß√µes baseadas no checklist de adapta√ß√£o
    instrucoes_checklist = ""
    if checklist_adaptacao and isinstance(checklist_adaptacao, dict):
        necessidades_ativas = []
        
        if checklist_adaptacao.get("questoes_desafiadoras"):
            necessidades_ativas.append("Aumentar o n√≠vel de desafio das quest√µes")
        else:
            necessidades_ativas.append("Manter ou reduzir o n√≠vel de dificuldade")
        
        if not checklist_adaptacao.get("compreende_instrucoes_complexas"):
            necessidades_ativas.append("Simplificar instru√ß√µes complexas")
        
        if checklist_adaptacao.get("instrucoes_passo_a_passo"):
            necessidades_ativas.append("Fornecer instru√ß√µes passo a passo")
        
        if checklist_adaptacao.get("dividir_em_etapas"):
            necessidades_ativas.append("Dividir quest√µes em etapas menores e mais gerenci√°veis")
        
        if checklist_adaptacao.get("paragrafos_curtos"):
            necessidades_ativas.append("Usar par√°grafos curtos para melhor compreens√£o")
        
        if checklist_adaptacao.get("dicas_apoio"):
            necessidades_ativas.append("Incluir dicas de apoio para resolver quest√µes")
        
        if not checklist_adaptacao.get("compreende_figuras_linguagem"):
            necessidades_ativas.append("Reduzir ou eliminar figuras de linguagem e infer√™ncias")
        
        if checklist_adaptacao.get("descricao_imagens"):
            necessidades_ativas.append("Incluir descri√ß√£o detalhada de imagens")
        
        if necessidades_ativas:
            instrucoes_checklist = f"""
    
    CHECKLIST DE ADAPTA√á√ÉO (baseado no PEI):
    {chr(10).join([f"- {n}" for n in necessidades_ativas])}
    
    REGRA CR√çTICA: N√ÉO aplique todas as necessidades em uma √∫nica quest√£o. 
    Para cada quest√£o, escolha APENAS 1-2 necessidades que fa√ßam mais sentido no contexto.
    Analise cada quest√£o individualmente e selecione as adapta√ß√µes mais relevantes.
    """
    
    prompt = f"""
    ESPECIALISTA EM DUA E INCLUS√ÉO. {style}
    1. ANALISE O PERFIL: {aluno.get('ia_sugestao', '')[:1000]}
    2. ADAPTE A PROVA: Use o hiperfoco ({aluno.get('hiperfoco', 'Geral')}) em 30% das quest√µes.
    {instrucoes_checklist}
    
    REGRA ABSOLUTA DE IMAGENS: O professor indicou imagens nas quest√µes: {lista_q if lista_q else "nenhuma"}.
    MANTENHA AS IMAGENS NO MESMO LOCAL ONDE ESTAVAM NA PROVA ORIGINAL.
    Para quest√µes com imagens, a estrutura OBRIGAT√ìRIA √©: 1. Enunciado -> 2. [[IMG_n√∫mero]] -> 3. Alternativas.
    NUNCA remova ou mova imagens de sua posi√ß√£o original.
    
    SA√çDA OBRIGAT√ìRIA (Use EXATAMENTE este divisor):
    [AN√ÅLISE PEDAG√ìGICA]
    ...an√°lise...
    ---DIVISOR---
    [ATIVIDADE]
    ...atividade...
    
    CONTEXTO: {materia} | {tema}. {"REMOVA GABARITO." if remover_resp else ""}
    TEXTO ORIGINAL: {texto}
    """
    
    try:
        resp = client.chat.completions.create(
            model="gpt-4o-mini", 
            messages=[{"role": "user", "content": prompt}], 
            temperature=0.7 if modo_profundo else 0.4
        )
        full_text = resp.choices[0].message.content
        if "---DIVISOR---" in full_text:
            parts = full_text.split("---DIVISOR---")
            return parts[0].replace("[AN√ÅLISE PEDAG√ìGICA]", "").strip(), parts[1].replace("[ATIVIDADE]", "").strip()
        return "An√°lise indispon√≠vel.", full_text
    except Exception as e:
        return str(e), ""

def adaptar_conteudo_imagem(api_key, aluno, imagem_bytes, materia, tema, tipo_atv, livro_professor, modo_profundo=False, checklist_adaptacao=None, imagem_separada=None):
    """Adapta conte√∫do de uma imagem para o estudante"""
    client = OpenAI(api_key=api_key)
    if not imagem_bytes:
        return "Erro: Imagem vazia", ""
    
    b64 = base64.b64encode(imagem_bytes).decode('utf-8')
    instrucao_livro = "ATEN√á√ÉO: IMAGEM COM RESPOSTAS. Remova todo gabarito/respostas." if livro_professor else ""
    style = "Fa√ßa uma an√°lise cr√≠tica para melhor adapta√ß√£o." if modo_profundo else "Transcreva e adapte."
    
    # Buscar hiperfoco do aluno
    hiperfoco = aluno.get('hiperfoco', 'Geral') or 'Geral'
    
    # Instru√ß√£o sobre imagem separada
    instrucao_imagem_separada = ""
    if imagem_separada:
        instrucao_imagem_separada = "\n    - O professor recortou a imagem separadamente para melhor qualidade. Use a tag [[IMG_2]] para inserir esta imagem recortada no local apropriado da quest√£o adaptada."
    
    # Montar instru√ß√µes baseadas no checklist de adapta√ß√£o (espec√≠fico para quest√£o √∫nica)
    instrucoes_checklist = ""
    if checklist_adaptacao and isinstance(checklist_adaptacao, dict):
        necessidades_ativas = []
        
        if checklist_adaptacao.get("questoes_desafiadoras"):
            necessidades_ativas.append("Aumentar o n√≠vel de desafio da quest√£o")
        else:
            necessidades_ativas.append("Manter ou reduzir o n√≠vel de dificuldade")
        
        if not checklist_adaptacao.get("compreende_instrucoes_complexas"):
            necessidades_ativas.append("Simplificar instru√ß√µes complexas")
        
        if checklist_adaptacao.get("instrucoes_passo_a_passo"):
            necessidades_ativas.append("Fornecer instru√ß√µes passo a passo detalhadas")
        
        if checklist_adaptacao.get("dividir_em_etapas"):
            necessidades_ativas.append("Dividir a quest√£o em etapas menores e mais gerenci√°veis")
        
        if checklist_adaptacao.get("paragrafos_curtos"):
            necessidades_ativas.append("Usar par√°grafos curtos para melhor compreens√£o")
        
        if checklist_adaptacao.get("dicas_apoio"):
            necessidades_ativas.append("Incluir dicas de apoio espec√≠ficas para resolver esta quest√£o")
        
        if not checklist_adaptacao.get("compreende_figuras_linguagem"):
            necessidades_ativas.append("Reduzir ou eliminar figuras de linguagem e infer√™ncias")
        
        if checklist_adaptacao.get("descricao_imagens"):
            necessidades_ativas.append("Incluir descri√ß√£o detalhada da imagem presente na quest√£o")
        
        if necessidades_ativas:
            instrucoes_checklist = f"""
    
    CHECKLIST DE ADAPTA√á√ÉO (baseado no PEI - QUEST√ÉO √öNICA):
    {chr(10).join([f"- {n}" for n in necessidades_ativas])}
    
    REGRA CR√çTICA: Como esta √© uma quest√£o √∫nica, aplique as adapta√ß√µes selecionadas de forma espec√≠fica e pontual.
    Selecione as 2-3 adapta√ß√µes mais relevantes para esta quest√£o espec√≠fica e aplique-as de forma integrada.
    """
    
    prompt = f"""
    ATUAR COMO: Especialista em Acessibilidade e OCR. {style}
    1. Transcreva o texto da imagem. {instrucao_livro}
    2. Adapte para o estudante (PEI: {aluno.get('ia_sugestao', '')[:800]}).
    3. HIPERFOCO ({hiperfoco}): Use o hiperfoco do estudante sempre que poss√≠vel para conectar e engajar na quest√£o. 
       Se o hiperfoco for relevante ao conte√∫do, integre-o naturalmente na adapta√ß√£o.
    {instrucoes_checklist}
    4. REGRA ABSOLUTA DE IMAGEM: 
    - Se a quest√£o original tinha imagem, detecte-a na imagem fornecida e insira a tag [[IMG_1]] no mesmo local onde estava.
    {instrucao_imagem_separada}
    - MANTENHA AS IMAGENS NO MESMO LOCAL ONDE ESTAVAM NA QUEST√ÉO ORIGINAL.
    
    SA√çDA OBRIGAT√ìRIA (Respeite o divisor):
    [AN√ÅLISE PEDAG√ìGICA]
    ...an√°lise...
    ---DIVISOR---
    [ATIVIDADE]
    ...atividade...
    """
    
    # Preparar mensagens com imagem(s)
    content_msgs = [
        {"type": "text", "text": prompt}, 
        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64}"}}
    ]
    
    # Se houver imagem separada, adicionar tamb√©m
    if imagem_separada:
        b64_separada = base64.b64encode(imagem_separada).decode('utf-8')
        content_msgs.append({
            "type": "text", 
            "text": "IMAGEM RECORTADA SEPARADAMENTE PELO PROFESSOR (use tag [[IMG_2]] para inserir no local apropriado):"
        })
        content_msgs.append({
            "type": "image_url", 
            "image_url": {"url": f"data:image/jpeg;base64,{b64_separada}"}
        })
    
    msgs = [
        {
            "role": "user", 
            "content": content_msgs
        }
    ]
    
    try:
        resp = client.chat.completions.create(
            model="gpt-4o-mini", 
            messages=msgs, 
            temperature=0.7 if modo_profundo else 0.4
        )
        full_text = resp.choices[0].message.content
        analise = "An√°lise indispon√≠vel."
        atividade = full_text
        if "---DIVISOR---" in full_text:
            parts = full_text.split("---DIVISOR---")
            analise = parts[0].replace("[AN√ÅLISE PEDAG√ìGICA]", "").strip()
            atividade = parts[1].replace("[ATIVIDADE]", "").strip()
        atividade = garantir_tag_imagem(atividade)
        return analise, atividade
    except Exception as e:
        return str(e), ""

def criar_profissional(api_key, aluno, materia, objeto, qtd, tipo_q, qtd_imgs, verbos_bloom=None, habilidades_bncc=None, modo_profundo=False, checklist_adaptacao=None):
    """Cria atividade profissional do zero"""
    client = OpenAI(api_key=api_key)
    hiperfoco = aluno.get('hiperfoco', 'Geral')
    
    # Instru√ß√£o de imagens
    instrucao_img = f"Incluir imagens em {qtd_imgs} quest√µes (use [[GEN_IMG: termo]]). REGRA DE POSI√á√ÉO: A tag da imagem ([[GEN_IMG: termo]]) DEVE vir logo AP√ìS o enunciado e ANTES das alternativas." if qtd_imgs > 0 else "Sem imagens."
    
    # Instru√ß√£o de Bloom
    instrucao_bloom = ""
    if verbos_bloom:
        lista_verbos = ", ".join(verbos_bloom)
        instrucao_bloom = f"""
        6. TAXONOMIA DE BLOOM (RIGOROSO):
           - Utilize OBRIGATORIAMENTE os seguintes verbos de a√ß√£o selecionados: {lista_verbos}.
           - Distribua esses verbos entre as quest√µes criadas.
           - REGRA DE FORMATA√á√ÉO: O verbo de comando deve vir no in√≠cio do enunciado, em **NEGRITO E CAIXA ALTA** (Ex: **ANALISE**, **IDENTIFIQUE**, **EXPLIQUE**).
           - Use apenas UM verbo de comando por quest√£o.
        """
    
    # Instru√ß√£o de habilidades BNCC
    instrucao_habilidades = ""
    if habilidades_bncc:
        habilidades_str = "\n".join([f"- {hab}" for hab in habilidades_bncc])
        instrucao_habilidades = f"""
        7. HABILIDADES BNCC (RIGOROSO):
           - Alinhe as quest√µes com as seguintes habilidades da BNCC:
           {habilidades_str}
           - Inclua refer√™ncias √†s habilidades nos enunciados quando pertinente.
        """
    
    # Montar instru√ß√µes do checklist de adapta√ß√£o (cria√ß√£o j√° considerando PEI)
    instrucoes_checklist = ""
    if checklist_adaptacao and isinstance(checklist_adaptacao, dict):
        necessidades_ativas = []
        if checklist_adaptacao.get("questoes_desafiadoras"):
            necessidades_ativas.append("Incluir quest√µes mais desafiadoras")
        else:
            necessidades_ativas.append("Manter n√≠vel de dificuldade acess√≠vel")
        if not checklist_adaptacao.get("compreende_instrucoes_complexas"):
            necessidades_ativas.append("Usar instru√ß√µes simples e diretas")
        if checklist_adaptacao.get("instrucoes_passo_a_passo"):
            necessidades_ativas.append("Fornecer instru√ß√µes passo a passo")
        if checklist_adaptacao.get("dividir_em_etapas"):
            necessidades_ativas.append("Dividir em etapas menores")
        if checklist_adaptacao.get("paragrafos_curtos"):
            necessidades_ativas.append("Usar par√°grafos curtos")
        if checklist_adaptacao.get("dicas_apoio"):
            necessidades_ativas.append("Incluir dicas de apoio")
        if not checklist_adaptacao.get("compreende_figuras_linguagem"):
            necessidades_ativas.append("Evitar figuras de linguagem complexas")
        if checklist_adaptacao.get("descricao_imagens"):
            necessidades_ativas.append("Incluir descri√ß√£o de imagens quando houver")
        if necessidades_ativas:
            instrucoes_checklist = f"""
    8. CHECKLIST DE ADAPTA√á√ÉO (baseado no PEI):
       {chr(10).join([f"- {n}" for n in necessidades_ativas])}
       Aplique essas orienta√ß√µes de forma coerente nas quest√µes criadas.
    """
    
    # Formato baseado no tipo
    if tipo_q == "Discursiva":
        diretriz_tipo = "3. FORMATO DISCURSIVO (RIGOROSO): Crie apenas quest√µes abertas. N√ÉO inclua alternativas, apenas linhas para resposta."
    else:
        diretriz_tipo = "3. FORMATO OBJETIVO: Crie quest√µes de m√∫ltipla escolha com distratores inteligentes."
    
    style = "Atue como uma banca examinadora rigorosa." if modo_profundo else "Atue como professor elaborador."
    
    prompt = f"""
    {style}
    Crie prova de {materia} ({objeto}). QTD: {qtd} ({tipo_q}).
    
    DIRETRIZES: 
    1. Contexto Real. 
    2. Hiperfoco ({hiperfoco}) em 30%. 
    {diretriz_tipo}
    4. Imagens: {instrucao_img} (NUNCA repita a mesma imagem). 
    5. Divis√£o Clara.
    
    REGRA DE OURO GRAMATICAL (IMPERATIVO):
    - TODOS os comandos das quest√µes devem estar no modo IMPERATIVO (Ex: "Cite", "Explique", "Calcule", "Analise", "Escreva").
    - JAMAIS use o infinitivo (Ex: "Citar", "Explicar", "Calcular").
    - Se houver verbos de Bloom selecionados, CONJUGUE-OS para o IMPERATIVO.
    - O verbo de comando deve vir no in√≠cio do enunciado, em **NEGRITO E CAIXA ALTA** (Ex: **ANALISE**, **IDENTIFIQUE**).
    
    {instrucao_bloom}
    {instrucao_habilidades}
    {instrucoes_checklist}
    
    SA√çDA OBRIGAT√ìRIA:
    [AN√ÅLISE PEDAG√ìGICA]
    ...an√°lise...
    ---DIVISOR---
    [ATIVIDADE]
    ...quest√µes...
    """
    
    try:
        resp = client.chat.completions.create(
            model="gpt-4o-mini", 
            messages=[{"role": "user", "content": prompt}], 
            temperature=0.8 if modo_profundo else 0.6
        )
        full_text = resp.choices[0].message.content
        if "---DIVISOR---" in full_text:
            parts = full_text.split("---DIVISOR---")
            return parts[0].replace("[AN√ÅLISE PEDAG√ìGICA]", "").strip(), parts[1].replace("[ATIVIDADE]", "").strip()
        return "An√°lise indispon√≠vel.", full_text
    except Exception as e:
        return str(e), ""

def gerar_experiencia_ei_bncc(api_key, aluno, campo_exp, objetivo, feedback_anterior=""):
    """Gera experi√™ncia para Educa√ß√£o Infantil"""
    client = OpenAI(api_key=api_key)
    hiperfoco = aluno.get('hiperfoco', 'Brincar')
    
    ajuste_prompt = ""
    if feedback_anterior:
        ajuste_prompt = f"AJUSTE SOLICITADO PELO PROFESSOR: {feedback_anterior}. Refa√ßa considerando isso."

    prompt = f"""
    ATUAR COMO: Especialista em Educa√ß√£o Infantil (BNCC) e Inclus√£o.
    ESTUDANTE: {aluno['nome']} (Educa√ß√£o Infantil).
    HIPERFOCO: {hiperfoco}.
    RESUMO DAS NECESSIDADES (PEI): {aluno.get('ia_sugestao', '')[:600]}
    
    SUA MISS√ÉO: Criar uma EXPERI√äNCIA L√öDICA, CONCRETA E VISUAL focada no Campo de Experi√™ncia: "{campo_exp}".
    Objetivo Espec√≠fico: {objetivo}
    {ajuste_prompt}
    
    REGRAS:
    1. N√£o crie "provas" ou "folhinhas". Crie VIV√äNCIAS.
    2. Use o hiperfoco para engajar (ex: se gosta de dinossauros, conte dinossauros).
    3. Liste materiais concretos (massinha, tinta, blocos).
    4. D√™ o passo a passo para o professor.
    
    SA√çDA ESPERADA (Markdown):
    ## üß∏ Experi√™ncia: [Nome Criativo]
    **üéØ Intencionalidade:** ...
    **üì¶ Materiais:** ...
    **üë£ Como Acontece:** ...
    **üé® Adapta√ß√£o para {aluno['nome'].split()[0]}:** ...
    """
    
    try:
        resp = client.chat.completions.create(
            model="gpt-4o-mini", 
            messages=[{"role": "user", "content": prompt}], 
            temperature=0.7
        )
        return resp.choices[0].message.content
    except Exception as e:
        return str(e)

def gerar_roteiro_aula_completo(api_key, aluno, materia, assunto, habilidades_bncc=None, verbos_bloom=None, ano=None, unidade_tematica=None, objeto_conhecimento=None, feedback_anterior=""):
    """Gera roteiro de aula completo com BNCC"""
    client = OpenAI(api_key=api_key)
    
    # Construir informa√ß√µes da BNCC
    info_bncc = ""
    if habilidades_bncc:
        info_bncc += f"\nHabilidades BNCC:"
        for hab in habilidades_bncc:
            info_bncc += f"\n- {hab}"
    
    if ano:
        info_bncc += f"\nAno: {ano}"
    
    if unidade_tematica:
        info_bncc += f"\nUnidade Tem√°tica: {unidade_tematica}"
    
    if objeto_conhecimento:
        info_bncc += f"\nObjeto do Conhecimento: {objeto_conhecimento}"
    
    # Construir informa√ß√µes Bloom
    info_bloom = ""
    if verbos_bloom:
        info_bloom = f"\nVerbos da Taxonomia de Bloom: {', '.join(verbos_bloom)}"
    
    ajuste = f"\nAJUSTES SOLICITADOS: {feedback_anterior}" if feedback_anterior else ""
    
    prompt = f"""
    Crie um ROTEIRO DE AULA INDIVIDUALIZADO para {aluno['nome']}.
    
    INFORMA√á√ïES DO ESTUDANTE:
    - Perfil: {aluno.get('ia_sugestao', '')[:500]}
    - Hiperfoco: {aluno.get('hiperfoco', 'Geral')}
    
    INFORMA√á√ïES DA AULA:
    - Componente Curricular: {materia}
    - Assunto: {assunto}
    {info_bncc}
    {info_bloom}
    {ajuste}
    
    ESTRUTURA OBRIGAT√ìRIA:
    
    1. **CONEX√ÉO INICIAL COM O HIPERFOCO** (2-3 minutos)
       - Como conectar o tema com o hiperfoco do estudante
    
    2. **OBJETIVOS DA AULA**
       - Objetivos claros e mensur√°veis
    
    3. **DESENVOLVIMENTO PASSO A PASSO** (15-20 minutos)
       - Divida em 3-4 etapas claras
       - Inclua perguntas mediadoras
       - Use exemplos relacionados ao hiperfoco
    
    4. **ATIVIDADE PR√ÅTICA INDIVIDUAL** (5-7 minutos)
       - Tarefa que o estudante pode fazer sozinho
    
    5. **FECHAMENTO E REFLEX√ÉO** (3-5 minutos)
       - Verifica√ß√£o dos objetivos
       - Pergunta de reflex√£o
    
    6. **RECURSOS E MATERIAIS**
    
    7. **AVALIA√á√ÉO FORMATIVA**
       - Como avaliar durante a aula
    """
    
    try:
        resp = client.chat.completions.create(
            model="gpt-4o-mini", 
            messages=[{"role": "user", "content": prompt}], 
            temperature=0.7
        )
        return resp.choices[0].message.content
    except Exception as e:
        return str(e)

def gerar_dinamica_inclusiva_completa(api_key, aluno, materia, assunto, qtd_alunos, caracteristicas_turma, habilidades_bncc=None, verbos_bloom=None, ano=None, unidade_tematica=None, objeto_conhecimento=None, feedback_anterior=""):
    """Gera din√¢mica inclusiva completa com BNCC"""
    client = OpenAI(api_key=api_key)
    
    # Construir informa√ß√µes da BNCC
    info_bncc = ""
    if habilidades_bncc:
        info_bncc += f"\nHabilidades BNCC:"
        for hab in habilidades_bncc:
            info_bncc += f"\n- {hab}"
    
    if ano:
        info_bncc += f"\nAno: {ano}"
    
    if unidade_tematica:
        info_bncc += f"\nUnidade Tem√°tica: {unidade_tematica}"
    
    if objeto_conhecimento:
        info_bncc += f"\nObjeto do Conhecimento: {objeto_conhecimento}"
    
    # Construir informa√ß√µes Bloom
    info_bloom = ""
    if verbos_bloom:
        info_bloom = f"\nVerbos da Taxonomia de Bloom: {', '.join(verbos_bloom)}"
    
    ajuste = f"\nAJUSTES SOLICITADOS: {feedback_anterior}" if feedback_anterior else ""
    
    prompt = f"""
    Crie uma DIN√ÇMICA INCLUSIVA para {qtd_alunos} estudantes.
    
    INFORMA√á√ïES DO ESTUDANTE FOCAL:
    - Nome: {aluno['nome']}
    - Perfil: {aluno.get('ia_sugestao', '')[:400]}
    - Hiperfoco: {aluno.get('hiperfoco', 'Geral')}
    
    INFORMA√á√ïES DA DIN√ÇMICA:
    - Componente Curricular: {materia}
    - Tema: {assunto}
    - Caracter√≠sticas da turma: {caracteristicas_turma}
    {info_bncc}
    {info_bloom}
    {ajuste}
    
    ESTRUTURA OBRIGAT√ìRIA:
    
    1. **NOME DA DIN√ÇMICA E OBJETIVO**
       - Nome criativo
       - Objetivo claro
    
    2. **MATERIAIS NECESS√ÅRIOS**
    
    3. **PREPARA√á√ÉO**
       - Como preparar a sala/ambiente
    
    4. **PASSO A PASSO** (detalhado)
       - Instru√ß√µes claras para o professor
       - Inclua adapta√ß√µes para o estudante focal
    
    5. **DURA√á√ÉO ESTIMADA**
    
    6. **AVALIA√á√ÉO**
       - Como avaliar a participa√ß√£o de todos
    
    7. **VARIA√á√ïES**
       - Sugest√µes para adaptar a din√¢mica
    """
    
    try:
        resp = client.chat.completions.create(
            model="gpt-4o-mini", 
            messages=[{"role": "user", "content": prompt}], 
            temperature=0.7
        )
        return resp.choices[0].message.content
    except Exception as e:
        return str(e)

def gerar_plano_aula_completo(api_key, materia, assunto, metodologia, tecnica, qtd_alunos, recursos, habilidades_bncc=None, verbos_bloom=None, ano=None, unidade_tematica=None, objeto_conhecimento=None, aluno_info=None):
    """Gera plano de aula completo com BNCC"""
    client = OpenAI(api_key=api_key)
    
    # Construir informa√ß√µes da BNCC
    info_bncc = ""
    if habilidades_bncc:
        info_bncc += f"\nHABILIDADES BNCC:"
        for hab in habilidades_bncc:
            info_bncc += f"\n- {hab}"
    
    if ano:
        info_bncc += f"\nAno: {ano}"
    
    if unidade_tematica:
        info_bncc += f"\nUnidade Tem√°tica: {unidade_tematica}"
    
    if objeto_conhecimento:
        info_bncc += f"\nObjeto do Conhecimento: {objeto_conhecimento}"
    
    # Construir informa√ß√µes Bloom
    info_bloom = ""
    if verbos_bloom:
        info_bloom = f"\nVERBOS DA TAXONOMIA DE BLOOM: {', '.join(verbos_bloom)}"
    
    # Informa√ß√µes do aluno (para DUA)
    info_aluno = ""
    if aluno_info:
        info_aluno = f"""
    INFORMA√á√ïES DO ESTUDANTE (DUA):
    - Nome: {aluno_info.get('nome', '')}
    - Hiperfoco: {aluno_info.get('hiperfoco', '')}
    - Perfil: {aluno_info.get('ia_sugestao', '')[:300]}
        """
    
    prompt = f"""
    ATUAR COMO: Coordenador Pedag√≥gico Especialista em BNCC, DUA e Metodologias Ativas.
    
    Crie um PLANO DE AULA COMPLETO com as seguintes informa√ß√µes:
    
    INFORMA√á√ïES B√ÅSICAS:
    - Componente Curricular: {materia}
    - Tema/Assunto: {assunto}
    - Metodologia: {metodologia}
    - T√©cnica: {tecnica if tecnica else 'N√£o especificada'}
    - Quantidade de Estudantes: {qtd_alunos}
    - Recursos Dispon√≠veis: {', '.join(recursos)}
    
    {info_bncc}
    {info_bloom}
    {info_aluno}
    
    ESTRUTURA DO PLANO (Markdown):
    
    ## üìã PLANO DE AULA: {assunto}
    
    ### üéØ OBJETIVOS DE APRENDIZAGEM
    - Objetivo geral
    - Objetivos espec√≠ficos (3-4)
    - Habilidades da BNCC trabalhadas
    
    ### üìö CONTE√öDOS
    - Conte√∫dos conceituais
    - Conte√∫dos procedimentais
    - Conte√∫dos atitudinais
    
    ### ‚è∞ TEMPO ESTIMADO
    - Dura√ß√£o total: __ minutos
    
    ### üõ† RECURSOS DID√ÅTICOS
    - Lista de recursos necess√°rios
    
    ### üöÄ DESENVOLVIMENTO DA AULA
    #### 1. ACOLHIDA E MOTIVA√á√ÉO (__ minutos)
    - Atividade de engajamento
    
    #### 2. APRESENTA√á√ÉO DO CONTE√öDO (__ minutos)
    - Explica√ß√£o do tema
    - Conex√µes com conhecimentos pr√©vios
    
    #### 3. ATIVIDADE PR√ÅTICA (__ minutos)
    - Descri√ß√£o detalhada da atividade
    - Papel do professor
    - Papel dos estudantes
    
    #### 4. SOCIALIZA√á√ÉO (__ minutos)
    - Compartilhamento dos resultados
    - Discuss√£o coletiva
    
    #### 5. AVALIA√á√ÉO (__ minutos)
    - Instrumentos de avalia√ß√£o
    - Crit√©rios
    
    ### ‚ôø ADAPTA√á√ïES DUA (DESIGN UNIVERSAL PARA APRENDIZAGEM)
    - Engajamento: Como manter todos motivados
    - Representa√ß√£o: M√∫ltiplas formas de apresentar o conte√∫do
    - A√ß√£o e Express√£o: M√∫ltiplas formas de expressar o aprendizado
    
    ### üìù AVALIA√á√ÉO
    - Avalia√ß√£o diagn√≥stica
    - Avalia√ß√£o formativa
    - Avalia√ß√£o somativa
    
    ### üîÑ RECUPERA√á√ÉO
    - Estrat√©gias para estudantes com dificuldades
    
    ### üìö REFER√äNCIAS
    - Refer√™ncias bibliogr√°ficas
    - Sites recomendados
    """
    
    try:
        resp = client.chat.completions.create(
            model="gpt-4o-mini", 
            messages=[{"role": "user", "content": prompt}], 
            temperature=0.7
        )
        return resp.choices[0].message.content
    except Exception as e:
        return str(e)

def gerar_quebra_gelo_profundo(api_key, aluno, materia, assunto, hiperfoco, tema_turma_extra=""):
    """Gera quebra-gelo profundo para engajamento"""
    client = OpenAI(api_key=api_key)
    
    prompt = f"""
    Crie 3 sugest√µes de 'Papo de Mestre' (Quebra-gelo/Introdu√ß√£o) para conectar o estudante {aluno['nome']} √† aula.
    Componente Curricular: {materia}. Assunto: {assunto}.
    Hiperfoco do estudante: {hiperfoco}.
    Tema de interesse da turma (DUA): {tema_turma_extra if tema_turma_extra else 'N√£o informado'}.
    
    O objetivo √© usar o hiperfoco ou o interesse da turma como UMA PONTE (estrat√©gia DUA de engajamento) para explicar o conceito de {assunto}.
    Seja criativo e profundo.
    """
    
    try:
        resp = client.chat.completions.create(
            model="gpt-4o-mini", 
            messages=[{"role": "user", "content": prompt}], 
            temperature=0.8
        )
        return resp.choices[0].message.content
    except Exception as e:
        return str(e)

    # ==============================================================================
# FUN√á√ïES DE BNCC (DROPDOWNS)
# ==============================================================================

def ano_celula_contem(ano_celula, ano_busca):
    """
    Verifica se ano_busca est√° na c√©lula Ano da BNCC.
    A c√©lula pode ter m√∫ltiplos anos: "1¬∫, 2¬∫, 3¬∫, 4¬∫, 5¬∫" (ex: Arte, L√≠ngua Portuguesa).
    Retorna True se ano_busca (ex: "3¬∫") est√° na lista.
    """
    if not ano_busca or not ano_celula:
        return False
    cell = str(ano_celula).strip()
    busca = str(ano_busca).strip()
    partes = [p.strip() for p in cell.split(",")]
    return busca in partes


def extrair_ano_bncc_do_aluno(aluno):
    """
    Extrai o ano/s√©rie no formato BNCC a partir do estudante (grade/serie do PEI).
    Retorna ex: "3¬∫", "7¬∫", "1EM" ou None.
    """
    if not aluno or not isinstance(aluno, dict):
        return None
    serie = aluno.get("serie") or aluno.get("grade") or ""
    if not serie or not isinstance(serie, str):
        return None
    s = str(serie).strip()
    # 1¬™ S√©rie (EM) -> 1EM
    m_em = re.search(r"(\d)\s*¬™?\s*s√©rie", s, re.IGNORECASE)
    if m_em or "em" in s.lower() or "m√©dio" in s.lower():
        n = m_em.group(1) if m_em else re.search(r"(\d)", s)
        if n:
            return f"{int(n.group(1))}EM"
    # 7¬∫ Ano, 3¬∫ ano, etc.
    m = re.search(r"(\d\s*¬∫)", s)
    return m.group(1).replace(" ", "") if m else None


def padronizar_ano(ano_str):
    """Converte diferentes formatos de ano para um padr√£o orden√°vel"""
    if not isinstance(ano_str, str):
        ano_str = str(ano_str)
    
    ano_str = ano_str.strip()
    
    # Padr√µes comuns
    padroes = [
        (r'(\d+)\s*¬∫?\s*ano', 'ano'),
        (r'(\d+)\s*¬™?\s*s√©rie', 'ano'),
        (r'(\d+)\s*em', 'em'),
        (r'ef\s*(\d+)', 'ano'),
        (r'(\d+)\s*per√≠odo', 'ano'),
        (r'(\d+)\s*semestre', 'ano'),
    ]
    
    for padrao, tipo in padroes:
        match = re.search(padrao, ano_str.lower())
        if match:
            num = match.group(1)
            if tipo == 'em':
                return f"{int(num):02d}EM"
            else:
                return f"{int(num):02d}"
    
    return ano_str

def ordenar_anos(anos_lista):
    """Ordena anos de forma inteligente"""
    anos_padronizados = []
    
    for ano in anos_lista:
        padrao = padronizar_ano(str(ano))
        anos_padronizados.append((padrao, ano))
    
    anos_padronizados.sort(key=lambda x: x[0])
    return [ano_original for _, ano_original in anos_padronizados]

@st.cache_data
def carregar_bncc_completa():
    """Carrega o CSV da BNCC com todas as colunas necess√°rias"""
    try:
        if not os.path.exists('bncc.csv'):
            st.warning("üìÑ Arquivo 'bncc.csv' n√£o encontrado na pasta do script")
            return None
        
        try:
            df = pd.read_csv('bncc.csv', delimiter=',', encoding='utf-8')
        except:
            try:
                df = pd.read_csv('bncc.csv', delimiter=';', encoding='utf-8')
            except Exception as e:
                st.error(f"‚ùå Erro ao ler CSV: {str(e)[:100]}")
                return None
        
        colunas_necessarias = ['Ano', 'Disciplina', 'Unidade Tem√°tica', 
                              'Objeto do Conhecimento', 'Habilidade']
        
        colunas_faltando = []
        for col in colunas_necessarias:
            if col not in df.columns:
                colunas_faltando.append(col)
        
        if colunas_faltando:
            st.error(f"‚ùå Colunas faltando: {colunas_faltando}")
            return None
        
        df = df.dropna(subset=['Ano', 'Disciplina', 'Objeto do Conhecimento'])
        df['Ano'] = df['Ano'].astype(str).str.strip()
        df['Disciplina'] = df['Disciplina'].str.replace('Ed. F√≠sica', 'Educa√ß√£o F√≠sica')
        
        return df
    
    except Exception as e:
        st.error(f"‚ùå Erro: {str(e)[:100]}")
        return None

def criar_dropdowns_bncc_completos_melhorado(key_suffix="", mostrar_habilidades=True, aluno=None):
    """
    Cria 5 dropdowns hier√°rquicos conectados com multiselect para habilidades.
    Se aluno for passado, pr√©-preenche o Ano com a s√©rie do estudante (PEI).
    """
    # Carregar dados se necess√°rio
    if 'bncc_df_completo' not in st.session_state:
        st.session_state.bncc_df_completo = carregar_bncc_completa()
    
    dados = st.session_state.bncc_df_completo
    
    # Se n√£o tem dados, mostrar campos b√°sicos
    if dados is None or dados.empty:
        st.warning("‚ö†Ô∏è BNCC n√£o carregada. Usando campos b√°sicos.")
        anos_opcoes = ordenar_anos(["1", "2", "3", "4", "5", "6", "7", "8", "9", "1EM", "2EM", "3EM"])
        ano_padrao = extrair_ano_bncc_do_aluno(aluno) if aluno else None
        col1, col2, col3 = st.columns(3)
        with col1:
            if ano_padrao:
                rotulo = f"{ano_padrao} ano" if "EM" not in str(ano_padrao).upper() else f"{str(ano_padrao).replace('EM','')}¬™ s√©rie EM"
                st.selectbox("Ano (PEI)", [rotulo], index=0, disabled=True, key=f"ano_basico_{key_suffix}")
                ano = ano_padrao
            else:
                ano = st.selectbox("Ano (PEI)", anos_opcoes, key=f"ano_basico_{key_suffix}")
        with col2:
            disc_opcoes, disc_mostrar, disc_default = _get_hub_componentes_filtrados()
            disciplina = st.selectbox(
                "Componente Curricular",
                disc_opcoes,
                index=0 if disc_opcoes else 0,
                disabled=not disc_mostrar,
                key=f"disc_basico_{key_suffix}"
            )
        with col3:
            objeto = st.text_input("Objeto do Conhecimento", placeholder="Ex: Fra√ß√µes", 
                                  key=f"obj_basico_{key_suffix}")
        
        col4, col5 = st.columns(2)
        with col4:
            unidade = st.text_input("Unidade Tem√°tica", placeholder="Ex: N√∫meros", 
                                   key=f"unid_basico_{key_suffix}")
        
        habilidades = []
        if mostrar_habilidades:
            with col5:
                habilidades_selecionadas = st.multiselect(
                    "Habilidades (selecione uma ou mais)",
                    ["Digite abaixo...", "Habilidade 1", "Habilidade 2"],
                    default=[],
                    key=f"hab_multi_basico_{key_suffix}"
                )
                
                if "Digite abaixo..." in habilidades_selecionadas:
                    habilidade_texto = st.text_area("Digite as habilidades:", 
                                                   placeholder="Uma por linha", 
                                                   key=f"hab_texto_basico_{key_suffix}")
                    habilidades = habilidade_texto.split('\n') if habilidade_texto else []
                else:
                    habilidades = habilidades_selecionadas
        
        return ano, disciplina, unidade, objeto, habilidades
    
    # TEMOS DADOS - Ano fixo pelo PEI (sem escolha), Componente filtrado por ano
    # O PEI √© o documento norteador: ano/s√©rie vem dele. O sistema determina, n√£o o usu√°rio.
    
    ano_selecionado = extrair_ano_bncc_do_aluno(aluno) if aluno else None
    if not ano_selecionado and aluno:
        # Fallback: tenta pegar do pei_data.serie
        pei = (aluno.get("pei_data") or {}) if isinstance(aluno.get("pei_data"), dict) else {}
        serie_pei = pei.get("serie") or aluno.get("grade") or ""
        ano_selecionado = extrair_ano_bncc_do_aluno({"serie": serie_pei, "grade": serie_pei}) if serie_pei else None
    
    def _rotulo_ano(a):
        """Formata ano para exibi√ß√£o: 7¬∫ -> 7¬∫ ano, 1EM -> 1¬™ s√©rie EM"""
        if not a:
            return ""
        s = str(a).strip()
        if "EM" in s.upper():
            n = s.replace("EM", "").strip()
            return f"{n}¬™ s√©rie EM" if n else s
        return f"{s} ano"
    
    if not ano_selecionado:
        st.warning("‚ö†Ô∏è **S√©rie n√£o informada no PEI.** Informe o ano/s√©rie do estudante na aba Estudante do PEI para filtrar a BNCC.")
        anos_originais = dados['Ano'].dropna().unique().tolist()
        anos_ordenados = ordenar_anos(anos_originais)
        ano_selecionado = st.selectbox("Ano (escolha enquanto o PEI n√£o tiver s√©rie)", anos_ordenados, key=f"ano_bncc_fallback_{key_suffix}")
    
    # Linha 1: Ano (desabilitado, do PEI), Componente Curricular, Unidade Tem√°tica
    col1, col2, col3 = st.columns(3)
    
    with col1:
        rotulo = _rotulo_ano(ano_selecionado)
        st.selectbox("Ano (PEI)", [rotulo], index=0, disabled=True, key=f"ano_bncc_display_{key_suffix}")
    
    with col2:
        # Componente Curricular ‚Äî apenas disciplinas que t√™m conte√∫do para o ano do PEI; filtrado por componente do professor
        if ano_selecionado:
            mask_ano = dados['Ano'].apply(lambda x: str(x).strip() == str(ano_selecionado).strip() or ano_celula_contem(x, ano_selecionado))
            df_por_ano = dados[mask_ano]
            disciplinas_raw = sorted(df_por_ano['Disciplina'].dropna().unique())
            disc_opcoes, disc_mostrar, disc_default = _filtrar_disciplinas_por_membro(disciplinas_raw)
            disciplina_selecionada = st.selectbox(
                "Componente Curricular",
                disc_opcoes,
                index=0,
                disabled=not disc_mostrar,
                key=f"disc_bncc_{key_suffix}"
            )
        else:
            disciplina_selecionada = None
            mask_ano = None
    
    with col3:
        if ano_selecionado and disciplina_selecionada and mask_ano is not None:
            unid_filtradas = dados[mask_ano & (dados['Disciplina'] == disciplina_selecionada)]
            unidades = sorted(unid_filtradas['Unidade Tem√°tica'].dropna().unique())
            unidade_selecionada = st.selectbox("Unidade Tem√°tica", unidades, key=f"unid_bncc_{key_suffix}")
        else:
            unidade_selecionada = None
    
    # Linha 2: Objeto do Conhecimento
    st.markdown("---")
    col4 = st.columns(1)[0]
    obj_filtrados = pd.DataFrame()
    
    with col4:
        if ano_selecionado and disciplina_selecionada and unidade_selecionada:
            df_filtrado = dados[mask_ano & (dados['Disciplina'] == disciplina_selecionada)]
            obj_filtrados = df_filtrado[df_filtrado['Unidade Tem√°tica'] == unidade_selecionada]
            objetos = sorted(obj_filtrados['Objeto do Conhecimento'].dropna().unique())
            
            if objetos:
                objeto_selecionado = st.selectbox("Objeto do Conhecimento", objetos, key=f"obj_bncc_{key_suffix}")
            else:
                objeto_selecionado = st.text_input("Objeto do Conhecimento", 
                                                  placeholder="N√£o encontrado, digite", 
                                                  key=f"obj_input_bncc_{key_suffix}")
        else:
            objeto_selecionado = st.text_input("Objeto do Conhecimento", 
                                              placeholder="Selecione primeiro", 
                                              key=f"obj_wait_bncc_{key_suffix}")
    
    # Habilidades (opcional)
    habilidades_selecionadas = []
    if mostrar_habilidades:
        st.markdown("---")
        col5 = st.columns(1)[0]
        
        with col5:
            if (ano_selecionado and disciplina_selecionada and 
                unidade_selecionada and objeto_selecionado and 
                isinstance(objeto_selecionado, str) and not objeto_selecionado.startswith("Selecione")):
                
                hab_filtradas = obj_filtrados[obj_filtrados['Objeto do Conhecimento'] == objeto_selecionado] if not obj_filtrados.empty else obj_filtrados
                todas_habilidades = sorted(hab_filtradas['Habilidade'].dropna().unique())
                
                if todas_habilidades:
                    st.markdown(f"**üîç {len(todas_habilidades)} habilidade(s) encontrada(s):**")
                    
                    opcoes_habilidades = st.multiselect(
                        "Selecione uma ou mais habilidades:",
                        todas_habilidades,
                        default=todas_habilidades[:min(3, len(todas_habilidades))],
                        key=f"hab_multi_bncc_{key_suffix}"
                    )
                    
                    with st.expander("‚ûï Adicionar habilidade personalizada"):
                        habilidade_extra = st.text_area(
                            "Digite habilidades adicionais (uma por linha):",
                            placeholder="Ex:\nEF05MA01 - Ler n√∫meros\nEF05MA02 - Comparar n√∫meros",
                            key=f"hab_extra_bncc_{key_suffix}"
                        )
                        
                        if habilidade_extra:
                            habilidades_extras = [h.strip() for h in habilidade_extra.split('\n') if h.strip()]
                            opcoes_habilidades.extend(habilidades_extras)
                    
                    habilidades_selecionadas = opcoes_habilidades
                else:
                    st.info("‚ÑπÔ∏è Nenhuma habilidade encontrada para este objeto.")
                    habilidades_padrao = st.multiselect(
                        "Selecione ou adicione habilidades:",
                        ["Digite abaixo..."],
                        default=[],
                        key=f"hab_vazio_bncc_{key_suffix}"
                    )
                    
                    if "Digite abaixo..." in habilidades_padrao:
                        habilidade_texto = st.text_area("Digite as habilidades:", 
                                                       placeholder="Uma por linha", 
                                                       key=f"hab_texto_{key_suffix}")
                        habilidades_selecionadas = habilidade_texto.split('\n') if habilidade_texto else []
                    else:
                        habilidades_selecionadas = habilidades_padrao
            else:
                st.info("‚ÑπÔ∏è Selecione Componente Curricular, Unidade e Objeto para ver as habilidades.")
                habilidades_selecionadas = []
    
    return (ano_selecionado, disciplina_selecionada, unidade_selecionada, 
            objeto_selecionado, habilidades_selecionadas)

def criar_dropdowns_bncc_simplificado(key_suffix=""):
    """Cria dropdowns simplificados da BNCC (apenas at√© objeto do conhecimento)"""
    return criar_dropdowns_bncc_completos_melhorado(key_suffix=key_suffix, mostrar_habilidades=False)

    # ==============================================================================
# COMPONENTES DE INTERFACE (UI)
# ==============================================================================

def aplicar_estilos():
    """Aplica estilos CSS espec√≠ficos do Hub (sem hero card - j√° renderizado acima)"""
    st.markdown("""
    <style>
        /* PEDAGOGIA BOX */
        .pedagogia-box { 
            background-color: #F8FAFC; 
            border-left: 4px solid #CBD5E1; 
            padding: 20px; 
            border-radius: 0 12px 12px 0; 
            margin-bottom: 25px; 
            font-size: 0.95rem; 
            color: #4A5568; 
        }
    
        /* RESOURCE BOX */
        .resource-box { 
            background: #F8FAFC; 
            border: 1px solid #E2E8F0; 
            border-radius: 12px; 
            padding: 20px; 
            margin: 15px 0; 
        }
        
        /* ACTION BUTTONS */
        .action-buttons { 
            display: flex; 
            gap: 10px; 
            margin-top: 20px; 
            flex-wrap: wrap; 
        }
        
        /* TIMELINE STYLES */
        .timeline-header { 
            background: white; 
            border-radius: 12px; 
            padding: 20px;
            margin-bottom: 20px; 
            border: 1px solid #E2E8F0;
            display: flex; 
            align-items: center; 
            justify-content: space-between; 
        }
        .prog-bar-bg { 
            width: 100%; 
            height: 8px; 
            background: #E2E8F0; 
            border-radius: 4px; 
            overflow: hidden; 
            margin-top: 8px; 
        }
        .prog-bar-fill { 
            height: 100%; 
            background: linear-gradient(90deg, #06B6D4, #0891B2); 
            transition: width 1s; 
        }
        
        /* HEADER DO ESTUDANTE */
        .student-header {
            background-color: #F8FAFC;
            border: 1px solid #E2E8F0;
            border-radius: 16px;
            padding: 18px 24px;
            margin-bottom: 18px;
            display: flex; 
            justify-content: space-between; 
            align-items: center;
            box-shadow: 0 2px 4px rgba(0,0,0,0.02);
        }
        .student-info-item {
            display: flex;
            flex-direction: column;
            gap: 4px;
        }
        .student-label {
            font-size: 0.78rem; 
            color: #64748B; 
            font-weight: 800;
            text-transform: uppercase; 
            letter-spacing: 1px;
        }
        .student-value { 
            font-size: 1.15rem; 
            color: #1E293B; 
            font-weight: 800; 
        }
        
        /* RESPONSIVIDADE */
        @media (max-width: 768px) { 
            .student-header { 
                flex-direction: column; 
                align-items:flex-start; 
                gap: 12px; 
            }
        }
    </style>
    """, unsafe_allow_html=True)

def render_cabecalho_aluno(aluno):
    """Renderiza o cabe√ßalho com informa√ß√µes do aluno"""
    # Garantir que o hiperfoco n√£o seja igual ao diagn√≥stico
    hiperfoco_display = aluno.get('hiperfoco', '-') or '-'
    diagnostico_aluno = aluno.get('diagnosis', '') or ''
    
    # Se o hiperfoco for igual ao diagn√≥stico, mostrar padr√£o
    if hiperfoco_display == diagnostico_aluno and hiperfoco_display != '-':
        hiperfoco_display = "Interesses gerais (A descobrir)"
    
    st.markdown(f"""
        <div class="student-header">
            <div class="student-info-item"><div class="student-label">Nome</div><div class="student-value">{aluno.get('nome')}</div></div>
            <div class="student-info-item"><div class="student-label">S√©rie</div><div class="student-value">{aluno.get('serie', '-')}</div></div>
            <div class="student-info-item"><div class="student-label">Hiperfoco</div><div class="student-value">{hiperfoco_display}</div></div>
        </div>
    """, unsafe_allow_html=True)

def criar_seletor_bloom(chave_prefixo):
    """Componente para Taxonomia de Bloom (Ajuste M√≠nimo)"""
    # Removido o get_icon de dentro do r√≥tulo para evitar a quebra
    usar_bloom = st.checkbox("Usar Taxonomia de Bloom (Revisada)", key=f"usar_bloom_{chave_prefixo}")
    
    mem_key = f"bloom_memoria_{chave_prefixo}"
    if mem_key not in st.session_state:
        st.session_state[mem_key] = {cat: [] for cat in TAXONOMIA_BLOOM.keys()}
    
    verbos_finais = []
    
    if usar_bloom:
        col_b1, col_b2 = st.columns(2)
        with col_b1:
            cat_atual = st.selectbox("Categoria Cognitiva:", list(TAXONOMIA_BLOOM.keys()), key=f"cat_bloom_{chave_prefixo}")
        with col_b2:
            # R√≥tulo limpo apenas com texto
            selecao_atual = st.multiselect(
                f"Verbos de '{cat_atual}':", 
                TAXONOMIA_BLOOM[cat_atual],
                default=st.session_state[mem_key].get(cat_atual, []),
                key=f"ms_bloom_{chave_prefixo}_{cat_atual}"
            )
            st.session_state[mem_key][cat_atual] = selecao_atual
        
        for cat in st.session_state[mem_key]:
            verbos_finais.extend(st.session_state[mem_key][cat])
            
    return verbos_finais if usar_bloom else None
    
# ==============================================================================
# FUN√á√ïES DAS ABAS PRINCIPAIS
# ==============================================================================

def render_aba_adaptar_prova(aluno, api_key):
    """Renderiza a aba de adapta√ß√£o de prova"""
    st.markdown("""
    <div class="pedagogia-box">
        <div class="pedagogia-title"><i class="ri-file-edit-line"></i> Adapta√ß√£o Curricular (DUA)</div>
        Transforme provas padr√£o em avalia√ß√µes acess√≠veis. O sistema simplifica enunciados, 
        insere imagens de apoio e ajusta o layout para reduzir a carga cognitiva.
    </div>
    """, unsafe_allow_html=True)
    
    # BNCC + Assunto em expander compacto
    with st.expander("üìö BNCC e Assunto", expanded=True):
        if 'bncc_df_completo' not in st.session_state:
            st.session_state.bncc_df_completo = carregar_bncc_completa()
        dados = st.session_state.bncc_df_completo
        
        # Ano determinado pelo PEI ‚Äî campo desabilitado para dar clareza (7¬∫ ano, 5¬∫ ano etc.)
        ano_bncc = extrair_ano_bncc_do_aluno(aluno) if aluno else None
        if not ano_bncc and aluno:
            pei = (aluno.get("pei_data") or {}) if isinstance(aluno.get("pei_data"), dict) else {}
            serie_pei = pei.get("serie") or aluno.get("grade") or ""
            ano_bncc = extrair_ano_bncc_do_aluno({"serie": serie_pei, "grade": serie_pei}) if serie_pei else None
        if not ano_bncc and dados is not None and not dados.empty:
            st.warning("‚ö†Ô∏è S√©rie n√£o informada no PEI. Informe na aba Estudante do PEI.")
            anos_ord = ordenar_anos(dados['Ano'].dropna().unique().tolist())
            ano_bncc = st.selectbox("Ano (fallback)", anos_ord, key="ano_adaptar_prova_fallback")
        # Linha 1: Ano (desabilitado), Componente, Unidade
        col_bncc1, col_bncc2, col_bncc3 = st.columns(3)
        with col_bncc1:
            if ano_bncc:
                rotulo = f"{ano_bncc} ano" if "EM" not in str(ano_bncc).upper() else f"{str(ano_bncc).replace('EM','')}¬™ s√©rie EM"
                st.selectbox("Ano (PEI)", [rotulo], index=0, disabled=True, key="ano_adaptar_prova_display")
            else:
                st.empty()
        with col_bncc2:
            if dados is not None and not dados.empty and ano_bncc:
                mask_ano = dados['Ano'].apply(lambda x: str(x).strip() == str(ano_bncc).strip() or ano_celula_contem(x, ano_bncc))
                df_por_ano = dados[mask_ano]
                disciplinas_raw = sorted(df_por_ano['Disciplina'].dropna().unique())
                disc_opcoes, disc_mostrar, _ = _filtrar_disciplinas_por_membro(disciplinas_raw)
                disciplina_bncc = st.selectbox("Componente Curricular", disc_opcoes, index=0, disabled=not disc_mostrar, key="disc_adaptar_prova_compact")
            else:
                lista_fallback = ["L√≠ngua Portuguesa", "Matem√°tica", "Ci√™ncias", "Hist√≥ria", "Geografia", "Arte", "Educa√ß√£o F√≠sica", "Ingl√™s"]
                disc_opcoes, disc_mostrar, _ = _filtrar_disciplinas_por_membro(lista_fallback)
                disciplina_bncc = st.selectbox("Componente Curricular", disc_opcoes, index=0, disabled=not disc_mostrar, key="disc_adaptar_prova_compact")
        with col_bncc3:
            if dados is not None and not dados.empty and ano_bncc and disciplina_bncc:
                mask_ano = dados['Ano'].apply(lambda x: str(x).strip() == str(ano_bncc).strip() or ano_celula_contem(x, ano_bncc))
                unid_filtradas = dados[mask_ano & (dados['Disciplina'] == disciplina_bncc)]
                unidades = sorted(unid_filtradas['Unidade Tem√°tica'].dropna().unique())
                if unidades:
                    unidade_bncc = st.selectbox("Unidade Tem√°tica", unidades, key="unid_adaptar_prova_compact")
                else:
                    unidade_bncc = st.text_input("Unidade Tem√°tica", placeholder="Ex: N√∫meros", 
                                               key="unid_adaptar_prova_compact")
            else:
                unidade_bncc = st.text_input("Unidade Tem√°tica", placeholder="Ex: N√∫meros", 
                                           key="unid_adaptar_prova_compact")
        
        # Linha 2: Objeto do Conhecimento e Assunto
        col_obj, col_ass = st.columns(2)
        with col_obj:
            if dados is not None and not dados.empty and ano_bncc and disciplina_bncc and unidade_bncc:
                mask_ano = dados['Ano'].apply(lambda x: str(x).strip() == str(ano_bncc).strip() or ano_celula_contem(x, ano_bncc))
                obj_filtrados = dados[mask_ano & (dados['Disciplina'] == disciplina_bncc) & (dados['Unidade Tem√°tica'] == unidade_bncc)]
                objetos = sorted(obj_filtrados['Objeto do Conhecimento'].dropna().unique())
                if objetos:
                    objeto_bncc = st.selectbox("Objeto do Conhecimento", objetos, key="obj_adaptar_prova_compact")
                else:
                    objeto_bncc = st.text_input("Objeto do Conhecimento", placeholder="Ex: Fra√ß√µes", 
                                               key="obj_adaptar_prova_compact")
            else:
                objeto_bncc = st.text_input("Objeto do Conhecimento", placeholder="Ex: Fra√ß√µes", 
                                           key="obj_adaptar_prova_compact")
        with col_ass:
            assunto_livre = st.text_input(
                "üìù Assunto (opcional)",
                value="",
                placeholder="Ex: Fra√ß√µes, Sistema Solar...",
                help="Preencha se quiser direcionar a adapta√ß√£o para um assunto espec√≠fico.",
                key="assunto_adaptar_prova_compact"
            )
    
    # Configura√ß√£o (Tipo e Upload)
    st.markdown("---")
    c1, c2 = st.columns([1, 2])
    tipo_d = c1.selectbox("Tipo de Documento", ["Prova", "Tarefa", "Avalia√ß√£o"], key="dtp")
    arquivo_d = c2.file_uploader("Upload do Arquivo DOCX", type=["docx"], key="fd")
    
    # Defini√ß√£o autom√°tica baseada na BNCC
    materia_d = disciplina_bncc if disciplina_bncc else "Geral"
    tema_d = objeto_bncc if objeto_bncc else "Geral"
    
    # Se assunto livre foi preenchido, usa ele como tema
    if assunto_livre and assunto_livre.strip():
        tema_d = assunto_livre.strip()
    
    if 'docx_imgs' not in st.session_state:
        st.session_state.docx_imgs = []
    if 'docx_txt' not in st.session_state:
        st.session_state.docx_txt = None
    
    if arquivo_d and arquivo_d.file_id != st.session_state.get('last_d'):
        st.session_state.last_d = arquivo_d.file_id
        txt, imgs = extrair_dados_docx(arquivo_d)
        st.session_state.docx_txt = txt
        st.session_state.docx_imgs = imgs
        st.success(f"{len(imgs)} imagens encontradas.")

    map_d = {}
    qs_d = []
    if st.session_state.docx_imgs:
        st.write("### Mapeamento de Imagens")
        cols = st.columns(3)
        for i, img in enumerate(st.session_state.docx_imgs):
            with cols[i % 3]:
                st.image(img, width=80)
                q = st.number_input(f"Pertence √† Quest√£o:", 0, 50, key=f"dq_{i}")
                if q > 0:
                    map_d[int(q)] = img
                    qs_d.append(int(q))

    # Inicializar checklist
    checklist_respostas = {}
    
    # Checklist de adapta√ß√£o - em expander retr√°til com checkboxes (4 colunas)
    with st.expander("üéØ Checklist de Adapta√ß√£o (baseado no PEI)", expanded=False):
        st.info("""
        ‚ö†Ô∏è **IMPORTANTE:** Marque apenas as adapta√ß√µes que devem ser aplicadas. 
        A IA escolher√° pontualmente 1-2 necessidades por quest√£o, evitando sobrecarga.
        """)
        
        # Organizar em 4 colunas para otimizar espa√ßo
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            precisa_desafio = st.checkbox(
                "Quest√µes mais desafiadoras",
                value=False,
                key="check_desafio"
            )
            precisa_passo = st.checkbox(
                "Instru√ß√µes passo a passo",
                value=False,
                key="check_passo"
            )
        
        with col2:
            compreende_complexas = st.checkbox(
                "Compreende instru√ß√µes complexas",
                value=False,
                key="check_complexas"
            )
            precisa_etapas = st.checkbox(
                "Dividir em etapas menores",
                value=False,
                key="check_etapas"
            )
        
        with col3:
            precisa_paragrafos_curtos = st.checkbox(
                "Par√°grafos curtos",
                value=False,
                key="check_paragrafos"
            )
            precisa_dicas = st.checkbox(
                "Dicas de apoio",
                value=False,
                key="check_dicas"
            )
        
        with col4:
            compreende_figuras = st.checkbox(
                "Compreende figuras de linguagem",
                value=False,
                key="check_figuras"
            )
            precisa_descricao_img = st.checkbox(
                "Descri√ß√£o de imagens",
                value=False,
                key="check_descricao"
            )
        
        # Compilar respostas em um dicion√°rio
        checklist_respostas = {
            "questoes_desafiadoras": precisa_desafio,
            "compreende_instrucoes_complexas": compreende_complexas,
            "instrucoes_passo_a_passo": precisa_passo,
            "dividir_em_etapas": precisa_etapas,
            "paragrafos_curtos": precisa_paragrafos_curtos,
            "dicas_apoio": precisa_dicas,
            "compreende_figuras_linguagem": compreende_figuras,
            "descricao_imagens": precisa_descricao_img
        }
    
    st.markdown("---")

    if st.button("üöÄ ADAPTAR PROVA", type="primary", key="btn_d", use_container_width=True):
        if not st.session_state.docx_txt:
            st.warning("Por favor, fa√ßa o upload de um arquivo DOCX.")
            st.stop()
        
        # Valida√ß√£o se BNCC foi preenchida (ou assunto livre)
        if not disciplina_bncc:
             st.warning("‚ö†Ô∏è Por favor, selecione a Disciplina nos campos da BNCC acima.")
             st.stop()
        
        if not objeto_bncc and not (assunto_livre and assunto_livre.strip()):
             st.warning("‚ö†Ô∏è Por favor, selecione o Objeto do Conhecimento (BNCC) ou preencha o campo Assunto acima para guiar a adapta√ß√£o.")
             st.stop()

        with st.spinner("A IA est√° analisando e adaptando o conte√∫do..."):
            # Salvar checklist no session_state para uso no refazer
            st.session_state['checklist_adaptacao_prova'] = checklist_respostas
            rac, txt = adaptar_conteudo_docx(
                api_key, aluno, st.session_state.docx_txt, materia_d, tema_d, tipo_d, True, qs_d,
                checklist_adaptacao=checklist_respostas
            )
            st.session_state['res_docx'] = {'rac': rac, 'txt': txt, 'map': map_d, 'valid': False}
            st.rerun()

    if 'res_docx' in st.session_state:
        res = st.session_state['res_docx']
        if res.get('valid'):
            st.success(f"{get_icon_emoji('validar')} **ATIVIDADE VALIDADA E PRONTA PARA USO**")
        else:
            col_v, col_r = st.columns([1, 1])
            if col_v.button(f"{get_icon_emoji('validar')} Validar", key="val_d", use_container_width=True):
                st.session_state['res_docx']['valid'] = True
                ou.track_ai_feedback("hub", "validated", content_type="prova_adaptada")
                st.rerun()
            if col_r.button(f"{get_icon_emoji('configurar')} Refazer (+Profundo)", key="redo_d", use_container_width=True):
                ou.track_ai_feedback("hub", "refazer", content_type="prova_adaptada")
                with st.spinner("Refazendo com an√°lise mais profunda..."):
                    # Recuperar checklist do session_state se dispon√≠vel
                    checklist_redo = st.session_state.get('checklist_adaptacao_prova', {})
                    rac, txt = adaptar_conteudo_docx(
                        api_key, aluno, st.session_state.docx_txt, materia_d, tema_d, tipo_d, True, qs_d, 
                        modo_profundo=True, checklist_adaptacao=checklist_redo
                    )
                    st.session_state['res_docx'] = {'rac': rac, 'txt': txt, 'map': map_d, 'valid': False}
                    st.rerun()

        st.markdown(f"<div class='analise-box'><div class='analise-title'>{get_icon('pedagogia', 20, '#06B6D4')} An√°lise Pedag√≥gica</div>{res['rac']}</div>", unsafe_allow_html=True)
        with st.container(border=True):
            partes = re.split(r'(\[\[IMG.*?\d+\]\])', res['txt'], flags=re.IGNORECASE)
            for p in partes:
                if "IMG" in p.upper() and re.search(r'\d+', p):
                    num = int(re.search(r'\d+', p).group(0))
                    im = res['map'].get(num)
                    if im:
                        st.image(im, width=300)
                elif p.strip():
                    st.markdown(p.strip())
        
        c_down1, c_down2 = st.columns(2)
        # Recuperar checklist para formata√ß√£o do Word
        checklist_formatacao = st.session_state.get('checklist_adaptacao_prova', {})
        docx = construir_docx_final(res['txt'], aluno, materia_d, res['map'], tipo_d, sem_cabecalho=False, checklist_adaptacao=checklist_formatacao)
        c_down1.download_button(
            f"{get_icon_emoji('download')} BAIXAR DOCX (Edit√°vel)", 
            docx, 
            "Prova_Adaptada.docx", 
            "primary",
            use_container_width=True
        )
        
        pdf_bytes = criar_pdf_generico(res['txt'])
        c_down2.download_button(
            f"{get_icon_emoji('download')} BAIXAR PDF (Visualiza√ß√£o)", 
            pdf_bytes, 
            "Prova_Adaptada.pdf", 
            mime="application/pdf", 
            type="secondary",
            use_container_width=True
        )

def render_aba_adaptar_atividade(aluno, api_key):
    """Renderiza a aba de adapta√ß√£o de atividade"""
    st.markdown("""
    <div class="pedagogia-box">
        <div class="pedagogia-title"><i class="ri-scissors-cut-line"></i> OCR & Adapta√ß√£o Visual</div>
        Tire foto de uma atividade do livro ou caderno. A IA extrai o texto, 
        remove polui√ß√£o visual e reestrutura o conte√∫do para o n√≠vel do aluno.
    </div>
    """, unsafe_allow_html=True)
    
    # BNCC + Assunto em expander compacto
    with st.expander("üìö BNCC e Assunto", expanded=True):
        if 'bncc_df_completo' not in st.session_state:
            st.session_state.bncc_df_completo = carregar_bncc_completa()
        dados = st.session_state.bncc_df_completo
        
        # Ano determinado pelo PEI ‚Äî campo desabilitado para clareza (7¬∫ ano, 5¬∫ ano etc.)
        ano_bncc = extrair_ano_bncc_do_aluno(aluno) if aluno else None
        if not ano_bncc and aluno:
            pei = (aluno.get("pei_data") or {}) if isinstance(aluno.get("pei_data"), dict) else {}
            serie_pei = pei.get("serie") or aluno.get("grade") or ""
            ano_bncc = extrair_ano_bncc_do_aluno({"serie": serie_pei, "grade": serie_pei}) if serie_pei else None
        if not ano_bncc and dados is not None and not dados.empty:
            st.warning("‚ö†Ô∏è S√©rie n√£o informada no PEI. Informe na aba Estudante do PEI.")
            anos_ord = ordenar_anos(dados['Ano'].dropna().unique().tolist())
            ano_bncc = st.selectbox("Ano (fallback)", anos_ord, key="ano_adaptar_atividade_fallback")
        # Linha 1: Ano (desabilitado), Componente, Unidade
        col_bncc1, col_bncc2, col_bncc3 = st.columns(3)
        with col_bncc1:
            if ano_bncc:
                rotulo = f"{ano_bncc} ano" if "EM" not in str(ano_bncc).upper() else f"{str(ano_bncc).replace('EM','')}¬™ s√©rie EM"
                st.selectbox("Ano (PEI)", [rotulo], index=0, disabled=True, key="ano_adaptar_atividade_display")
            else:
                st.empty()
        with col_bncc2:
            if dados is not None and not dados.empty and ano_bncc:
                mask_ano = dados['Ano'].apply(lambda x: str(x).strip() == str(ano_bncc).strip() or ano_celula_contem(x, ano_bncc))
                df_por_ano = dados[mask_ano]
                disciplinas_raw = sorted(df_por_ano['Disciplina'].dropna().unique())
                disc_opcoes, disc_mostrar, _ = _filtrar_disciplinas_por_membro(disciplinas_raw)
                disciplina_bncc = st.selectbox("Componente Curricular", disc_opcoes, index=0, disabled=not disc_mostrar, key="disc_adaptar_atividade_compact")
            else:
                lista_fallback = ["L√≠ngua Portuguesa", "Matem√°tica", "Ci√™ncias", "Hist√≥ria", "Geografia", "Arte", "Educa√ß√£o F√≠sica", "Ingl√™s"]
                disc_opcoes, disc_mostrar, _ = _filtrar_disciplinas_por_membro(lista_fallback)
                disciplina_bncc = st.selectbox("Componente Curricular", disc_opcoes, index=0, disabled=not disc_mostrar, key="disc_adaptar_atividade_compact")
        with col_bncc3:
            if dados is not None and not dados.empty and ano_bncc and disciplina_bncc:
                mask_ano = dados['Ano'].apply(lambda x: str(x).strip() == str(ano_bncc).strip() or ano_celula_contem(x, ano_bncc))
                unid_filtradas = dados[mask_ano & (dados['Disciplina'] == disciplina_bncc)]
                unidades = sorted(unid_filtradas['Unidade Tem√°tica'].dropna().unique())
                if unidades:
                    unidade_bncc = st.selectbox("Unidade Tem√°tica", unidades, key="unid_adaptar_atividade_compact")
                else:
                    unidade_bncc = st.text_input("Unidade Tem√°tica", placeholder="Ex: N√∫meros", 
                                               key="unid_adaptar_atividade_compact")
            else:
                unidade_bncc = st.text_input("Unidade Tem√°tica", placeholder="Ex: N√∫meros", 
                                           key="unid_adaptar_atividade_compact")
        
        # Linha 2: Objeto do Conhecimento e Assunto
        col_obj, col_ass = st.columns(2)
        with col_obj:
            if dados is not None and not dados.empty and ano_bncc and disciplina_bncc and unidade_bncc:
                mask_ano = dados['Ano'].apply(lambda x: str(x).strip() == str(ano_bncc).strip() or ano_celula_contem(x, ano_bncc))
                obj_filtrados = dados[mask_ano & (dados['Disciplina'] == disciplina_bncc) & (dados['Unidade Tem√°tica'] == unidade_bncc)]
                objetos = sorted(obj_filtrados['Objeto do Conhecimento'].dropna().unique())
                if objetos:
                    objeto_bncc = st.selectbox("Objeto do Conhecimento", objetos, key="obj_adaptar_atividade_compact")
                else:
                    objeto_bncc = st.text_input("Objeto do Conhecimento", placeholder="Ex: Fra√ß√µes", 
                                               key="obj_adaptar_atividade_compact")
            else:
                objeto_bncc = st.text_input("Objeto do Conhecimento", placeholder="Ex: Fra√ß√µes", 
                                           key="obj_adaptar_atividade_compact")
        with col_ass:
            assunto_livre = st.text_input(
                "üìù Assunto (opcional)",
                value="",
                placeholder="Ex: Fra√ß√µes, Sistema Solar...",
                help="Preencha se quiser direcionar a adapta√ß√£o para um assunto espec√≠fico.",
                key="assunto_adaptar_atividade_compact"
            )
    
    # Configura√ß√£o (Tipo e Upload)
    st.markdown("---")
    c1, c2 = st.columns([1, 2])
    tipo_i = c1.selectbox("Tipo", ["Atividade", "Tarefa", "Exerc√≠cio"], key="itp")
    arquivo_i = c2.file_uploader("Upload da Imagem/Foto", type=["png","jpg","jpeg"], key="fi")
    livro_prof = st.checkbox("üìï √â foto do Livro do Professor? (A IA remover√° as respostas)", value=False)
    
    # Defini√ß√£o autom√°tica baseada na BNCC
    materia_i = disciplina_bncc if disciplina_bncc else "Geral"
    tema_i = objeto_bncc if objeto_bncc else "Geral"
    
    # Se assunto livre foi preenchido, usa ele como tema
    if assunto_livre and assunto_livre.strip():
        tema_i = assunto_livre.strip()
    
    if 'img_raw' not in st.session_state:
        st.session_state.img_raw = None
    
    if arquivo_i and arquivo_i.file_id != st.session_state.get('last_i'):
        st.session_state.last_i = arquivo_i.file_id
        st.session_state.img_raw = sanitizar_imagem(arquivo_i.getvalue())

    # Processo de recorte
    cropped_res = None
    imagem_recortada = None
    
    if st.session_state.img_raw:
        st.markdown(f"### {icon_title('Passo 1: Recortar a Quest√£o', 'adaptar_atividade', 20, '#06B6D4')}", unsafe_allow_html=True)
        st.info("üí° **Importante:** Recorte a √°rea da quest√£o completa.")
        
        img_pil = Image.open(BytesIO(st.session_state.img_raw))
        img_pil.thumbnail((800, 800))
        cropped_res = st_cropper(img_pil, realtime_update=True, box_color='#FF0000', aspect_ratio=None, key="crop_i")
        if cropped_res:
            st.image(cropped_res, width=200, caption="Quest√£o recortada")
            
            # Se a quest√£o foi recortada, oferecer op√ß√£o de recortar imagem separadamente
            st.markdown("### üñºÔ∏è Passo 2: Recortar Imagem (Opcional)")
            st.caption("Se a quest√£o tem imagem e voc√™ quer recort√°-la separadamente para melhor qualidade, clique abaixo.")
            
            tem_imagem_na_questao = st.checkbox(
                "A quest√£o tem imagem e quero recort√°-la separadamente",
                value=False,
                key="tem_imagem_separada"
            )
            
            if tem_imagem_na_questao:
                st.info("üí° Recorte apenas a √°rea da imagem na quest√£o. Esta imagem ser√° inserida na quest√£o adaptada.")
                # Abrir a imagem original novamente para recortar s√≥ a imagem
                img_pil_original = Image.open(BytesIO(st.session_state.img_raw))
                img_pil_original.thumbnail((800, 800))
                imagem_recortada = st_cropper(img_pil_original, realtime_update=True, box_color='#00FF00', aspect_ratio=None, key="crop_img_separada")
                if imagem_recortada:
                    st.image(imagem_recortada, width=200, caption="Imagem recortada separadamente")

    # Inicializar checklist
    checklist_respostas = {}
    
    # Checklist de adapta√ß√£o - em expander retr√°til com checkboxes (4 colunas)
    with st.expander("üéØ Checklist de Adapta√ß√£o (baseado no PEI) - Quest√£o √önica", expanded=False):
        st.info("""
        ‚ö†Ô∏è **IMPORTANTE:** Marque apenas as adapta√ß√µes que devem ser aplicadas nesta quest√£o. 
        A IA aplicar√° as adapta√ß√µes selecionadas de forma pontual e espec√≠fica.
        """)
        
        # Organizar em 4 colunas para otimizar espa√ßo
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            precisa_desafio = st.checkbox(
                "Quest√µes mais desafiadoras",
                value=False,
                key="check_desafio_atividade"
            )
            precisa_passo = st.checkbox(
                "Instru√ß√µes passo a passo",
                value=False,
                key="check_passo_atividade"
            )
        
        with col2:
            compreende_complexas = st.checkbox(
                "Compreende instru√ß√µes complexas",
                value=False,
                key="check_complexas_atividade"
            )
            precisa_etapas = st.checkbox(
                "Dividir em etapas menores",
                value=False,
                key="check_etapas_atividade"
            )
        
        with col3:
            precisa_paragrafos_curtos = st.checkbox(
                "Par√°grafos curtos",
                value=False,
                key="check_paragrafos_atividade"
            )
            precisa_dicas = st.checkbox(
                "Dicas de apoio",
                value=False,
                key="check_dicas_atividade"
            )
        
        with col4:
            compreende_figuras = st.checkbox(
                "Compreende figuras de linguagem",
                value=False,
                key="check_figuras_atividade"
            )
            precisa_descricao_img = st.checkbox(
                "Descri√ß√£o de imagens",
                value=False,
                key="check_descricao_atividade"
            )
        
        # Compilar respostas em um dicion√°rio
        checklist_respostas = {
            "questoes_desafiadoras": precisa_desafio,
            "compreende_instrucoes_complexas": compreende_complexas,
            "instrucoes_passo_a_passo": precisa_passo,
            "dividir_em_etapas": precisa_etapas,
            "paragrafos_curtos": precisa_paragrafos_curtos,
            "dicas_apoio": precisa_dicas,
            "compreende_figuras_linguagem": compreende_figuras,
            "descricao_imagens": precisa_descricao_img
        }
        # Salvar no session_state para uso posterior
        st.session_state['checklist_adaptacao_atividade'] = checklist_respostas
    
    # Garantir que checklist_respostas est√° definido
    if not checklist_respostas:
        checklist_respostas = st.session_state.get('checklist_adaptacao_atividade', {})

    st.markdown("---")

    if st.button("üöÄ ADAPTAR ATIVIDADE", type="primary", key="btn_i", use_container_width=True):
        if not st.session_state.img_raw:
            st.warning("Por favor, envie uma imagem.")
            st.stop()
        
        # Valida√ß√£o BNCC (ou assunto livre)
        if not disciplina_bncc:
             st.warning("‚ö†Ô∏è Por favor, selecione a Disciplina nos campos da BNCC acima.")
             st.stop()
        
        if not objeto_bncc and not (assunto_livre and assunto_livre.strip()):
             st.warning("‚ö†Ô∏è Por favor, selecione o Objeto do Conhecimento (BNCC) ou preencha o campo Assunto acima para guiar a adapta√ß√£o.")
             st.stop()

        if not cropped_res:
            st.warning("‚ö†Ô∏è Por favor, recorte a √°rea da quest√£o primeiro.")
            st.stop()
        
        with st.spinner("Lendo imagem e adaptando conte√∫do..."):
            buf_c = BytesIO()
            cropped_res.convert('RGB').save(buf_c, format="JPEG", quality=90)
            img_bytes = buf_c.getvalue()
            
            # Processar imagem recortada separadamente se houver
            img_separada_bytes = None
            if imagem_recortada:
                buf_img = BytesIO()
                imagem_recortada.convert('RGB').save(buf_img, format="JPEG", quality=90)
                img_separada_bytes = buf_img.getvalue()
            
            # Salvar checklist no session_state para uso no refazer
            st.session_state['checklist_adaptacao_atividade'] = checklist_respostas
            st.session_state['img_separada_atividade'] = img_separada_bytes
            
            rac, txt = adaptar_conteudo_imagem(
                api_key, aluno, img_bytes, materia_i, tema_i, tipo_i, livro_prof, 
                checklist_adaptacao=checklist_respostas,
                imagem_separada=img_separada_bytes
            )
            
            # Mapear imagens: 1 = quest√£o recortada, 2 = imagem recortada separadamente (se houver)
            mapa_imgs = {1: img_bytes}
            if img_separada_bytes:
                mapa_imgs[2] = img_separada_bytes
            
            st.session_state['res_img'] = {'rac': rac, 'txt': txt, 'map': mapa_imgs, 'valid': False}
            st.rerun()

    if 'res_img' in st.session_state:
        res = st.session_state['res_img']
        if res.get('valid'):
            st.success(f"{get_icon_emoji('validar')} **ATIVIDADE VALIDADA E PRONTA PARA USO**")
        else:
            col_v, col_r = st.columns([1, 1])
            if col_v.button(f"{get_icon_emoji('validar')} Validar", key="val_i", use_container_width=True):
                st.session_state['res_img']['valid'] = True
                ou.track_ai_feedback("hub", "validated", content_type="atividade_adaptada")
                st.rerun()
            if col_r.button(f"{get_icon_emoji('configurar')} Refazer (+Profundo)", key="redo_i", use_container_width=True):
                ou.track_ai_feedback("hub", "refazer", content_type="atividade_adaptada")
                with st.spinner("Refazendo..."):
                    img_bytes = res['map'][1]
                    img_separada_redo = res['map'].get(2)  # Recuperar imagem separada se houver
                    # Recuperar checklist do session_state se dispon√≠vel
                    checklist_redo = st.session_state.get('checklist_adaptacao_atividade', {})
                    rac, txt = adaptar_conteudo_imagem(
                        api_key, aluno, img_bytes, materia_i, tema_i, tipo_i, livro_prof, 
                        modo_profundo=True, checklist_adaptacao=checklist_redo, imagem_separada=img_separada_redo
                    )
                    st.session_state['res_img'] = {'rac': rac, 'txt': txt, 'map': res['map'], 'valid': False}
                    st.rerun()

        st.markdown(f"<div class='analise-box'><div class='analise-title'>{get_icon('pedagogia', 20, '#06B6D4')} An√°lise Pedag√≥gica</div>{res['rac']}</div>", unsafe_allow_html=True)
        with st.container(border=True):
            partes = re.split(r'(\[\[IMG.*?\]\])', res['txt'], flags=re.IGNORECASE)
            for p in partes:
                if "IMG" in p.upper():
                    im = res['map'].get(1)
                    if im:
                        st.image(im, width=300)
                elif p.strip():
                    st.markdown(p.strip())
        
        c_down1, c_down2 = st.columns(2)
        # Recuperar checklist para formata√ß√£o do Word
        checklist_formatacao = st.session_state.get('checklist_adaptacao_atividade', {})
        docx = construir_docx_final(res['txt'], aluno, materia_i, res['map'], tipo_i, sem_cabecalho=False, checklist_adaptacao=checklist_formatacao)
        c_down1.download_button("üì• BAIXAR DOCX (Edit√°vel)", docx, "Atividade.docx", "primary", use_container_width=True)
        
        pdf_bytes = criar_pdf_generico(res['txt'])
        c_down2.download_button(
            f"{get_icon_emoji('download')} BAIXAR PDF (Visualiza√ß√£o)", 
            pdf_bytes, 
            "Atividade.pdf", 
            mime="application/pdf", 
            type="secondary",
            use_container_width=True
        )

def render_aba_criar_do_zero(aluno, api_key, unsplash_key):
    """Renderiza a aba de cria√ß√£o do zero"""
    st.markdown("""
    <div class="pedagogia-box">
        <div class="pedagogia-title"><i class="ri-magic-line"></i> Cria√ß√£o com DUA</div>
        Crie atividades do zero alinhadas ao PEI. A IA gera quest√µes contextualizadas, 
        usa o hiperfoco para engajamento. Imagens: prioridade no banco (Unsplash), IA s√≥ em √∫ltimo caso.
    </div>
    """, unsafe_allow_html=True)
    
    # BNCC + Assunto em expander compacto
    with st.expander("üìö BNCC e Assunto", expanded=True):
        if 'bncc_df_completo' not in st.session_state:
            st.session_state.bncc_df_completo = carregar_bncc_completa()
        dados = st.session_state.bncc_df_completo
        
        # Ano determinado pelo PEI ‚Äî campo desabilitado para clareza (7¬∫ ano, 5¬∫ ano etc.)
        ano_bncc = extrair_ano_bncc_do_aluno(aluno) if aluno else None
        if not ano_bncc and aluno:
            pei = (aluno.get("pei_data") or {}) if isinstance(aluno.get("pei_data"), dict) else {}
            serie_pei = pei.get("serie") or aluno.get("grade") or ""
            ano_bncc = extrair_ano_bncc_do_aluno({"serie": serie_pei, "grade": serie_pei}) if serie_pei else None
        if not ano_bncc and dados is not None and not dados.empty:
            st.warning("‚ö†Ô∏è S√©rie n√£o informada no PEI. Informe na aba Estudante do PEI.")
            anos_ord = ordenar_anos(dados['Ano'].dropna().unique().tolist())
            ano_bncc = st.selectbox("Ano (fallback)", anos_ord, key="ano_criar_zero_fallback")
        # Linha 1: Ano (desabilitado), Componente, Unidade
        c1, c2, c3 = st.columns(3)
        with c1:
            if ano_bncc:
                rotulo = f"{ano_bncc} ano" if "EM" not in str(ano_bncc).upper() else f"{str(ano_bncc).replace('EM','')}¬™ s√©rie EM"
                st.selectbox("Ano (PEI)", [rotulo], index=0, disabled=True, key="ano_criar_zero_display")
            else:
                st.empty()
        with c2:
            if dados is not None and not dados.empty and ano_bncc:
                mask_ano = dados['Ano'].apply(lambda x: str(x).strip() == str(ano_bncc).strip() or ano_celula_contem(x, ano_bncc))
                df_por_ano = dados[mask_ano]
                disciplinas_raw = sorted(df_por_ano['Disciplina'].dropna().unique())
                disc_opcoes, disc_mostrar, _ = _filtrar_disciplinas_por_membro(disciplinas_raw)
                disciplina_bncc = st.selectbox("Componente Curricular", disc_opcoes, index=0, disabled=not disc_mostrar, key="disc_criar_zero")
            else:
                lista_fallback = ["L√≠ngua Portuguesa", "Matem√°tica", "Ci√™ncias", "Hist√≥ria", "Geografia", "Arte", "Educa√ß√£o F√≠sica", "Ingl√™s"]
                disc_opcoes, disc_mostrar, _ = _filtrar_disciplinas_por_membro(lista_fallback)
                disciplina_bncc = st.selectbox("Componente Curricular", disc_opcoes, index=0, disabled=not disc_mostrar, key="disc_criar_zero")
        with c3:
            if dados is not None and not dados.empty and ano_bncc and disciplina_bncc:
                mask_ano = dados['Ano'].apply(lambda x: str(x).strip() == str(ano_bncc).strip() or ano_celula_contem(x, ano_bncc))
                unid_f = dados[mask_ano & (dados['Disciplina'] == disciplina_bncc)]
                unidades = sorted(unid_f['Unidade Tem√°tica'].dropna().unique())
                unidade_bncc = st.selectbox("Unidade Tem√°tica", unidades, key="unid_criar_zero") if unidades else st.text_input("Unidade Tem√°tica", placeholder="Ex: N√∫meros", key="unid_criar_zero")
            else:
                unidade_bncc = st.text_input("Unidade Tem√°tica", placeholder="Ex: N√∫meros", key="unid_criar_zero")
        
        # Linha 2: Objeto do Conhecimento, Assunto (opcional)
        c4, c5 = st.columns(2)
        with c4:
            if dados is not None and not dados.empty and ano_bncc and disciplina_bncc and unidade_bncc:
                mask_ano = dados['Ano'].apply(lambda x: str(x).strip() == str(ano_bncc).strip() or ano_celula_contem(x, ano_bncc))
                obj_f = dados[mask_ano & (dados['Disciplina'] == disciplina_bncc) & (dados['Unidade Tem√°tica'] == unidade_bncc)]
                objetos = sorted(obj_f['Objeto do Conhecimento'].dropna().unique())
                objeto_bncc = st.selectbox("Objeto do Conhecimento", objetos, key="obj_criar_zero") if objetos else st.text_input("Objeto do Conhecimento", placeholder="Ex: Fra√ß√µes", key="obj_criar_zero")
            else:
                objeto_bncc = st.text_input("Objeto do Conhecimento", placeholder="Ex: Fra√ß√µes", key="obj_criar_zero")
        with c5:
            assunto_criar = st.text_input("üìù Assunto (opcional)", value="", placeholder="Ex: Fra√ß√µes, Sistema Solar...", key="assunto_criar_zero", help="Direciona melhor o tema da atividade.")
        
        # Habilidades (compacto)
        habilidades_bncc = []
        if dados is not None and not dados.empty and ano_bncc and disciplina_bncc and unidade_bncc and objeto_bncc:
            mask_ano = dados['Ano'].apply(lambda x: str(x).strip() == str(ano_bncc).strip() or ano_celula_contem(x, ano_bncc))
            hab_f = dados[mask_ano & (dados['Disciplina'] == disciplina_bncc) & (dados['Unidade Tem√°tica'] == unidade_bncc) & (dados['Objeto do Conhecimento'] == objeto_bncc)]
            todas_hab = sorted(hab_f['Habilidade'].dropna().unique())
            if todas_hab:
                habilidades_bncc = st.multiselect("Habilidades BNCC", todas_hab, default=todas_hab[:min(3, len(todas_hab))], key="hab_criar_zero")
    
    mat_c = disciplina_bncc
    obj_c = assunto_criar.strip() if assunto_criar and assunto_criar.strip() else objeto_bncc
    
    # Configura√ß√£o da atividade (uma linha)
    st.markdown("---")
    r1, r2, r3, r4 = st.columns(4)
    with r1:
        qtd_c = st.slider("Quantidade de Quest√µes", 1, 10, 5, key="cq")
    with r2:
        tipo_quest = st.selectbox("Tipo de Quest√£o", ["Objetiva", "Discursiva"], key="ctq")
    with r3:
        usar_img = st.checkbox("Incluir Imagens?", value=True, key="usar_img")
    with r4:
        qtd_img_sel = st.slider("Qtd. imagens", 0, qtd_c, int(qtd_c/2) if qtd_c > 1 else 0, disabled=not usar_img, key="qtd_img_slider")
    
    # Inicializar vari√°veis antes dos expanders
    checklist_criar = {}
    verbos_finais_para_ia = []
    usar_bloom = False
    
    # Checklist e Taxonomia de Bloom juntos embaixo (lado a lado)
    col_check, col_bloom = st.columns(2)
    
    with col_check:
        with st.expander("üéØ Checklist de Adapta√ß√£o (opcional)", expanded=False):
            st.info("Marque as adapta√ß√µes que devem ser consideradas na cria√ß√£o da atividade.")
            col_c1, col_c2, col_c3, col_c4 = st.columns(4)
            with col_c1:
                check_desafio = st.checkbox("Quest√µes mais desafiadoras", value=False, key="c0_desafio")
                check_passo = st.checkbox("Instru√ß√µes passo a passo", value=False, key="c0_passo")
            with col_c2:
                check_complexas = st.checkbox("Compreende instru√ß√µes complexas", value=False, key="c0_complexas")
                check_etapas = st.checkbox("Dividir em etapas menores", value=False, key="c0_etapas")
            with col_c3:
                check_paragrafos = st.checkbox("Par√°grafos curtos", value=False, key="c0_paragrafos")
                check_dicas = st.checkbox("Dicas de apoio", value=False, key="c0_dicas")
            with col_c4:
                check_figuras = st.checkbox("Compreende figuras de linguagem", value=False, key="c0_figuras")
                check_descricao = st.checkbox("Descri√ß√£o de imagens", value=False, key="c0_descricao")
            checklist_criar = {
                "questoes_desafiadoras": check_desafio, "compreende_instrucoes_complexas": check_complexas,
                "instrucoes_passo_a_passo": check_passo, "dividir_em_etapas": check_etapas,
                "paragrafos_curtos": check_paragrafos, "dicas_apoio": check_dicas,
                "compreende_figuras_linguagem": check_figuras, "descricao_imagens": check_descricao
            }
    
    with col_bloom:
        with st.expander("üß† Taxonomia de Bloom (opcional)", expanded=False):
            if 'bloom_memoria' not in st.session_state:
                st.session_state.bloom_memoria = {cat: [] for cat in TAXONOMIA_BLOOM.keys()}
            usar_bloom = st.checkbox(f"{get_icon_emoji('configurar')} Usar Taxonomia de Bloom (Revisada)", key="usar_bloom")
            st.caption("[Saiba mais: Taxonomia de Bloom revisada (Wikipedia)](https://pt.wikipedia.org/wiki/Taxonomia_de_Bloom)")
            if usar_bloom:
                col_b1, col_b2 = st.columns(2)
                with col_b1:
                    cat_atual = st.selectbox("Categoria Cognitiva:", list(TAXONOMIA_BLOOM.keys()), key="cat_bloom")
                with col_b2:
                    selecao_atual = st.multiselect(f"Verbos de '{cat_atual}':", TAXONOMIA_BLOOM[cat_atual], default=st.session_state.bloom_memoria[cat_atual], key=f"ms_bloom_{cat_atual}")
                    st.session_state.bloom_memoria[cat_atual] = selecao_atual
                for cat in st.session_state.bloom_memoria:
                    verbos_finais_para_ia.extend(st.session_state.bloom_memoria[cat])
                if verbos_finais_para_ia:
                    st.info(f"**Verbos:** {', '.join(verbos_finais_para_ia)}")
    
    st.markdown("---")
    if st.button("‚ú® CRIAR ATIVIDADE", type="primary", key="btn_c", use_container_width=True):
        if not api_key:
            st.error(f"‚ùå Insira a chave da IA ({ou.AI_RED}) nas configura√ß√µes da sidebar.")
        else:
            with st.spinner("Elaborando atividade..."):
                qtd_final = qtd_img_sel if usar_img else 0
                rac, txt = criar_profissional(
                    api_key, aluno, mat_c, obj_c, qtd_c, tipo_quest, qtd_final,
                    verbos_bloom=verbos_finais_para_ia if usar_bloom else None,
                    habilidades_bncc=habilidades_bncc,
                    checklist_adaptacao=checklist_criar
                )
                novo_map = {}
                count = 0
                tags = re.findall(r'\[\[GEN_IMG: (.*?)\]\]', txt)
                for p in tags:
                    count += 1
                    # Prioridade: BANCO (Unsplash) primeiro; IA s√≥ em √∫ltimo caso
                    url = gerar_imagem_inteligente(api_key, p, unsplash_key, prioridade="BANCO")
                    if not url and unsplash_key:
                        url = gerar_imagem_inteligente(api_key, p, unsplash_key, prioridade="IA")
                    if url:
                        io = baixar_imagem_url(url)
                        if io:
                            novo_map[count] = io.getvalue()
                txt_fin = txt
                for i in range(1, count + 1):
                    txt_fin = re.sub(r'\[\[GEN_IMG: .*?\]\]', f"[[IMG_G{i}]]", txt_fin, count=1)
                st.session_state['res_create'] = {'rac': rac, 'txt': txt_fin, 'map': novo_map, 'valid': False, 'mat_c': mat_c, 'obj_c': obj_c, 'checklist': checklist_criar}
                st.rerun()
    
    # Exibi√ß√£o do resultado
    if 'res_create' in st.session_state:
        res = st.session_state['res_create']
        
        st.markdown("---")
        titulo_atividade = f"Atividade Criada: {res.get('mat_c', '')} - {res.get('obj_c', '')}"
        st.markdown(f"### {icon_title(titulo_atividade, 'criar_zero', 20, '#06B6D4')}", unsafe_allow_html=True)
        
        # Barra de status
        if res.get('valid'):
            st.success(f"{get_icon_emoji('validar')} **ATIVIDADE VALIDADA E PRONTA PARA USO**")
        else:
            col_val, col_ajust, col_desc = st.columns(3)
            with col_val:
                if st.button("‚úÖ Validar Atividade", key="val_c", use_container_width=True):
                    st.session_state['res_create']['valid'] = True
                    ou.track_ai_feedback("hub", "validated", content_type="atividade_criada")
                    st.rerun()
            with col_ajust:
                if st.button("üîÑ Refazer com Ajustes", key="redo_c", use_container_width=True):
                    ou.track_ai_feedback("hub", "refazer", content_type="atividade_criada")
                    st.session_state['res_create']['valid'] = False
                    st.info("Para ajustes, modifique os par√¢metros acima e clique em 'CRIAR ATIVIDADE' novamente.")
            with col_desc:
                if st.button("üóëÔ∏è Descartar", key="del_c", use_container_width=True):
                    del st.session_state['res_create']
                    st.rerun()
        
        # An√°lise Pedag√≥gica
        if res.get('rac'):
            with st.expander("üß† An√°lise Pedag√≥gica (clique para expandir)"):
                st.markdown(res['rac'])
        
        # Atividade Gerada
        st.markdown("#### üìù Atividade Gerada")
        with st.container(border=True):
            partes = re.split(r'(\[\[IMG_G\d+\]\])', res['txt'])
            for p in partes:
                tag = re.search(r'\[\[IMG_G(\d+)\]\]', p)
                if tag:
                    i = int(tag.group(1))
                    im = res['map'].get(i)
                    if im: 
                        st.image(im, width=300)
                elif p.strip(): 
                    st.markdown(p.strip())
        
        # Bot√µes de Download
        st.markdown("---")
        st.markdown("### üì• Download")
        col_down1, col_down2, col_down3 = st.columns(3)
        
        with col_down1:
            docx = construir_docx_final(res['txt'], aluno, mat_c, res.get('map', {}), "Criada", checklist_adaptacao=res.get('checklist'))
            st.download_button(
                label="üìÑ Baixar DOCX",
                data=docx,
                file_name=f"Atividade_{mat_c}_{date.today().strftime('%Y%m%d')}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        
        with col_down2:
            pdf_bytes = criar_pdf_generico(res['txt'])
            st.download_button(
                label="üìï Baixar PDF",
                data=pdf_bytes,
                file_name=f"Atividade_{mat_c}_{date.today().strftime('%Y%m%d')}.pdf",
                mime="application/pdf",
                use_container_width=True
            )
        
        with col_down3:
            st.download_button(
                label="üìù Baixar Texto",
                data=res['txt'],
                file_name=f"Atividade_{mat_c}_{date.today().strftime('%Y%m%d')}.txt",
                mime="text/plain",
                use_container_width=True
            )

def render_aba_estudio_visual(aluno, api_key, unsplash_key, gemini_key=None):
    """Renderiza a aba de est√∫dio visual (Gemini ou OpenAI + Unsplash)."""
    st.markdown(f"""
    <div class="pedagogia-box">
        <div class="pedagogia-title"><i class="ri-image-line"></i> Recursos Visuais</div>
        Gere flashcards, rotinas visuais e s√≠mbolos de comunica√ß√£o. Usa <strong>{ou.AI_BLUE}</strong> (imagens) ou <strong>{ou.AI_RED}</strong> (DALL-E)/Unsplash. Configure as chaves na sidebar.
    </div>
    """, unsafe_allow_html=True)
    
    col_scene, col_caa = st.columns(2)
    g_key = gemini_key or ou.get_gemini_api_key()
    
    with col_scene:
        st.markdown("#### üñºÔ∏è Ilustra√ß√£o")
        desc_m = st.text_area("Descreva a imagem:", height=100, key="vdm_padrao", placeholder="Ex: Sistema Solar simplificado com planetas coloridos...")
        hiperfoco_est = (aluno.get("hiperfoco") or "").strip()
        usar_hiperfoco_ilustracao = st.checkbox(
            "Usar hiperfoco do estudante como tema da ilustra√ß√£o",
            value=bool(hiperfoco_est),
            key="hub_usar_hiperfoco_ilustracao",
            help="Se marcado, abre o campo com o hiperfoco para voc√™ editar antes de gerar.",
        )
        tema_ilustracao = ""
        if usar_hiperfoco_ilustracao:
            tema_ilustracao = st.text_input(
                "Tema da ilustra√ß√£o (edite se quiser):",
                value=hiperfoco_est,
                key="hub_tema_ilustracao",
                placeholder="Ex: dinossauros, espa√ßo, m√∫sica...",
                help="Texto que ser√° usado como tema na gera√ß√£o. Vem do cadastro do estudante; voc√™ pode alterar.",
            )
            tema_ilustracao = (tema_ilustracao or "").strip()
        if st.button("üé® Gerar Imagem", key="btn_cena_padrao"):
            with st.spinner("Desenhando..."):
                prompt_completo = f"{desc_m}. Context: Education."
                if usar_hiperfoco_ilustracao and tema_ilustracao:
                    prompt_completo = f"Tema da ilustra√ß√£o: {tema_ilustracao}. {prompt_completo}"
                st.session_state.res_scene_url = gerar_imagem_inteligente(api_key, prompt_completo, unsplash_key, prioridade="IA", gemini_key=g_key)
                st.session_state.valid_scene = False

        if st.session_state.res_scene_url:
            st.image(st.session_state.res_scene_url)
            if st.session_state.valid_scene:
                st.success("Imagem validada!")
            else:
                c_vs1, c_vs2 = st.columns([1, 2])
                if c_vs1.button("‚úÖ Validar", key="val_sc_pd"):
                    st.session_state.valid_scene = True
                    ou.track_ai_feedback("hub", "validated", content_type="cena_ilustracao")
                    st.rerun()
                with c_vs2.expander("üîÑ Refazer Cena"):
                    fb_scene = st.text_input("Ajuste:", key="fb_sc_pd")
                    if st.button("Refazer", key="ref_sc_pd"):
                        ou.track_ai_feedback("hub", "refazer", content_type="cena_ilustracao", feedback_text=fb_scene or "")
                        with st.spinner("Redesenhando..."):
                            prompt_completo = f"{desc_m}. Context: Education."
                            if usar_hiperfoco_ilustracao and tema_ilustracao:
                                prompt_completo = f"Tema da ilustra√ß√£o: {tema_ilustracao}. {prompt_completo}"
                            st.session_state.res_scene_url = gerar_imagem_inteligente(api_key, prompt_completo, unsplash_key, feedback_anterior=fb_scene, prioridade="IA", gemini_key=g_key)
                            st.rerun()

    with col_caa:
        st.markdown("#### üó£Ô∏è S√≠mbolo CAA")
        palavra_chave = st.text_input("Conceito:", placeholder="Ex: Sil√™ncio", key="caa_input_padrao")
        
        if st.button("üß© Gerar Pictograma", key="btn_caa_padrao"):
            with st.spinner("Criando s√≠mbolo..."):
                st.session_state.res_caa_url = gerar_pictograma_caa(api_key, palavra_chave, gemini_key=g_key)
                st.session_state.valid_caa = False

        if st.session_state.res_caa_url:
            st.image(st.session_state.res_caa_url, width=300)
            if st.session_state.valid_caa:
                st.success("Pictograma validado!")
            else:
                c_vc1, c_vc2 = st.columns([1, 2])
                if c_vc1.button("‚úÖ Validar", key="val_caa_pd"):
                    st.session_state.valid_caa = True
                    ou.track_ai_feedback("hub", "validated", content_type="pictograma_caa")
                    st.rerun()
                with c_vc2.expander("üîÑ Refazer Picto"):
                    fb_caa = st.text_input("Ajuste:", key="fb_caa_pd")
                    if st.button("Refazer", key="ref_caa_pd"):
                        ou.track_ai_feedback("hub", "refazer", content_type="pictograma_caa", feedback_text=fb_caa or "")
                        with st.spinner("Recriando..."):
                            st.session_state.res_caa_url = gerar_pictograma_caa(api_key, palavra_chave, feedback_anterior=fb_caa, gemini_key=g_key)
                            st.rerun()

def render_aba_roteiro_individual(aluno, api_key):
    """Renderiza a aba de roteiro individual"""
    st.markdown("""
    <div class="pedagogia-box">
        <div class="pedagogia-title"><i class="ri-user-follow-line"></i> Roteiro de Aula Individualizado</div>
        Crie um passo a passo de aula <strong>espec√≠fico para este estudante</strong> do PEI. 
        A IA usar√° o hiperfoco como chave de acesso para o conte√∫do selecionado na BNCC.
    </div>
    """, unsafe_allow_html=True)
    
    # BNCC em expander
    with st.expander("üìö BNCC e Habilidades", expanded=True):
        ano_bncc, disciplina_bncc, unidade_bncc, objeto_bncc, habilidades_bncc = criar_dropdowns_bncc_completos_melhorado(key_suffix="roteiro", aluno=aluno)
    
    st.markdown("---")
    
    if st.button("üìù GERAR ROTEIRO INDIVIDUAL", type="primary", use_container_width=True):
        # Valida√ß√£o: Usa o objeto_bncc como assunto
        if objeto_bncc and habilidades_bncc:
            with st.spinner(f"Criando roteiro sobre '{objeto_bncc}'..."):
                res = gerar_roteiro_aula_completo(
                    api_key=api_key,
                    aluno=aluno,
                    materia=disciplina_bncc,
                    assunto=objeto_bncc, # Passa o objeto BNCC como assunto
                    habilidades_bncc=habilidades_bncc,
                    verbos_bloom=None, # Bloom removido
                    ano=ano_bncc,
                    unidade_tematica=unidade_bncc,
                    objeto_conhecimento=objeto_bncc
                )
                st.session_state['res_roteiro'] = res
                st.session_state['res_roteiro_valid'] = False
        else:
            st.warning("‚ö†Ô∏è Para gerar o roteiro, selecione a 'Disciplina', 'Objeto do Conhecimento' e pelo menos uma 'Habilidade' nos campos da BNCC acima.")
    
    # Exibi√ß√£o do resultado
    if 'res_roteiro' in st.session_state:
        res = st.session_state['res_roteiro']
        
        st.markdown("---")
        st.markdown("### üìã Roteiro de Aula Individualizado")
        
        if st.session_state.get('res_roteiro_valid'):
            st.success("‚úÖ **ROTEIRO VALIDADO E PRONTO PARA USO**")
        else:
            col_val, col_ajust, col_desc = st.columns(3)
            with col_val:
                if st.button("‚úÖ Validar Roteiro", key="val_roteiro", use_container_width=True):
                    st.session_state['res_roteiro_valid'] = True
                    ou.track_ai_feedback("hub", "validated", content_type="roteiro")
                    st.rerun()
            with col_ajust:
                if st.button("üîÑ Refazer com Ajustes", key="redo_roteiro", use_container_width=True):
                    ou.track_ai_feedback("hub", "refazer", content_type="roteiro")
                    st.session_state['res_roteiro_valid'] = False
                    st.info("Para ajustes, modifique as sele√ß√µes da BNCC e clique em 'GERAR ROTEIRO' novamente.")
            with col_desc:
                if st.button("üóëÔ∏è Descartar", key="del_roteiro", use_container_width=True):
                    del st.session_state['res_roteiro']
                    del st.session_state['res_roteiro_valid']
                    st.rerun()
        
        # Roteiro Gerado
        st.markdown(res)
        
        # Bot√µes de Download
        st.markdown("---")
        st.markdown("### üì• Download")
        col_down1, col_down2 = st.columns(2)
        
        with col_down1:
            docx_roteiro = criar_docx_simples(res, "Roteiro Individual")
            st.download_button(
                label="üìÑ Baixar DOCX",
                data=docx_roteiro,
                file_name=f"Roteiro_{disciplina_bncc}_{date.today().strftime('%Y%m%d')}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        
        with col_down2:
            pdf_bytes_roteiro = criar_pdf_generico(res)
            st.download_button(
                label="üìï Baixar PDF",
                data=pdf_bytes_roteiro,
                file_name=f"Roteiro_{disciplina_bncc}_{date.today().strftime('%Y%m%d')}.pdf",
                mime="application/pdf",
                use_container_width=True
            )

def render_aba_papo_mestre(aluno, api_key):
    """Renderiza a aba de papo de mestre"""
    st.markdown("""
    <div class="pedagogia-box">
        <div class="pedagogia-title"><i class="ri-chat-smile-2-line"></i> Engajamento & DUA (Papo de Mestre)</div>
        O hiperfoco √© um <strong>caminho neurol√≥gico</strong> j√° aberto. Use-o para conectar o estudante ao componente curricular.
        Aqui voc√™ tamb√©m pode adicionar um tema de interesse da turma toda (DUA) para criar conex√µes coletivas.
    </div>
    """, unsafe_allow_html=True)
    
    c1, c2 = st.columns(2)
    disc_opcoes, disc_mostrar, disc_default = _get_hub_componentes_filtrados()
    materia_papo = c1.selectbox(
        "Componente Curricular",
        disc_opcoes,
        index=0,
        disabled=not disc_mostrar,
        key="papo_mat"
    )
    assunto_papo = c2.text_input("Assunto da Aula:", key="papo_ass")
    
    c3, c4 = st.columns(2)
    hiperfoco_papo = c3.text_input("Hiperfoco (Estudante):", value=aluno.get('hiperfoco', 'Geral'), key="papo_hip")
    tema_turma = c4.text_input("Interesse da Turma (Opcional - DUA):", placeholder="Ex: Minecraft, Copa do Mundo...", key="papo_turma")
    
    if st.button("üó£Ô∏è CRIAR CONEX√ïES", type="primary"): 
        if assunto_papo:
            res = gerar_quebra_gelo_profundo(api_key, aluno, materia_papo, assunto_papo, hiperfoco_papo, tema_turma)
            st.session_state['res_papo'] = res
            st.session_state['res_papo_valid'] = False
        else:
            st.warning("Preencha o Assunto.")
    
    if 'res_papo' in st.session_state:
        res = st.session_state['res_papo']
        
        st.markdown("---")
        st.markdown("### üó£Ô∏è Conex√µes para Engajamento")
        
        if st.session_state.get('res_papo_valid'):
            st.success("‚úÖ **CONEX√ïES VALIDADAS E PRONTAS PARA USO**")
        else:
            col_val, col_ajust, col_desc = st.columns(3)
            with col_val:
                if st.button("‚úÖ Validar Conex√µes", key="val_papo", use_container_width=True):
                    st.session_state['res_papo_valid'] = True
                    ou.track_ai_feedback("hub", "validated", content_type="papo_mestre")
                    st.rerun()
            with col_ajust:
                if st.button("üîÑ Refazer com Ajustes", key="redo_papo", use_container_width=True):
                    ou.track_ai_feedback("hub", "refazer", content_type="papo_mestre")
                    st.session_state['res_papo_valid'] = False
                    st.info("Para ajustes, modifique os par√¢metros acima e clique em 'CRIAR CONEX√ïES' novamente.")
            with col_desc:
                if st.button("üóëÔ∏è Descartar", key="del_papo", use_container_width=True):
                    del st.session_state['res_papo']
                    del st.session_state['res_papo_valid']
                    st.rerun()
        
        st.markdown(res)
        
        # Bot√µes de Download
        st.markdown("---")
        st.markdown("### üì• Download")
        col_down1, col_down2 = st.columns(2)
        
        with col_down1:
            docx_papo = criar_docx_simples(res, "Papo de Mestre")
            st.download_button(
                label="üìÑ Baixar DOCX",
                data=docx_papo,
                file_name=f"Papo_Mestre_{date.today().strftime('%Y%m%d')}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        
        with col_down2:
            pdf_bytes_papo = criar_pdf_generico(res)
            st.download_button(
                label="üìï Baixar PDF",
                data=pdf_bytes_papo,
                file_name=f"Papo_Mestre_{date.today().strftime('%Y%m%d')}.pdf",
                mime="application/pdf",
                use_container_width=True
            )

def render_aba_dinamica_inclusiva(aluno, api_key):
    """Renderiza a aba de din√¢mica inclusiva"""
    st.markdown("""
    <div class="pedagogia-box">
        <div class="pedagogia-title"><i class="ri-group-line"></i> Din√¢mica Inclusiva</div>
        Atividades em grupo onde todos participam, respeitando as singularidades.
        Baseado no Objeto do Conhecimento selecionado.
    </div>
    """, unsafe_allow_html=True)
    
    # BNCC em expander
    with st.expander("üìö BNCC e Habilidades", expanded=True):
        ano_bncc, disciplina_bncc, unidade_bncc, objeto_bncc, habilidades_bncc = criar_dropdowns_bncc_completos_melhorado(key_suffix="dinamica", aluno=aluno)
    
    # Configura√ß√£o da Turma
    st.markdown("---")
    
    c3, c4 = st.columns(2)
    with c3:
        qtd_alunos = st.number_input("N√∫mero de Estudantes:", min_value=5, max_value=50, value=25, key="din_qtd")
    with c4:
        carac_turma = st.text_input(
            "Caracter√≠sticas da Turma (Opcional):", 
            placeholder="Ex: Turma agitada, gostam de competi√ß√£o...", 
            key="din_carac"
        )
    
    # Bot√£o para gerar
    st.markdown("---")
    
    if st.button("ü§ù CRIAR DIN√ÇMICA", type="primary", use_container_width=True): 
        if objeto_bncc and habilidades_bncc:
            with st.spinner(f"Criando din√¢mica sobre '{objeto_bncc}'..."):
                res = gerar_dinamica_inclusiva_completa(
                    api_key=api_key,
                    aluno=aluno,
                    materia=disciplina_bncc,
                    assunto=objeto_bncc, # Passa o objeto BNCC como assunto
                    qtd_alunos=qtd_alunos,
                    caracteristicas_turma=carac_turma,
                    habilidades_bncc=habilidades_bncc,
                    verbos_bloom=None, # Bloom Removido
                    ano=ano_bncc,
                    unidade_tematica=unidade_bncc,
                    objeto_conhecimento=objeto_bncc
                )
                st.session_state['res_dinamica'] = res
                st.session_state['res_dinamica_valid'] = False
        else:
            st.warning("‚ö†Ô∏è Selecione a 'Disciplina', 'Objeto do Conhecimento' e pelo menos uma 'Habilidade' nos campos da BNCC acima.")
    
    # Exibi√ß√£o do resultado
    if 'res_dinamica' in st.session_state:
        res = st.session_state['res_dinamica']
        
        st.markdown("---")
        st.markdown("### üìã Din√¢mica Inclusiva")
        
        if st.session_state.get('res_dinamica_valid'):
            st.success("‚úÖ **DIN√ÇMICA VALIDADA E PRONTA PARA USO**")
        else:
            col_val, col_ajust, col_desc = st.columns(3)
            with col_val:
                if st.button("‚úÖ Validar Din√¢mica", key="val_dinamica", use_container_width=True):
                    st.session_state['res_dinamica_valid'] = True
                    ou.track_ai_feedback("hub", "validated", content_type="dinamica")
                    st.rerun()
            with col_ajust:
                if st.button("üîÑ Refazer com Ajustes", key="redo_dinamica", use_container_width=True):
                    ou.track_ai_feedback("hub", "refazer", content_type="dinamica")
                    st.session_state['res_dinamica_valid'] = False
                    st.info("Para ajustes, modifique os par√¢metros e clique em 'CRIAR DIN√ÇMICA' novamente.")
            with col_desc:
                if st.button("üóëÔ∏è Descartar", key="del_dinamica", use_container_width=True):
                    del st.session_state['res_dinamica']
                    del st.session_state['res_dinamica_valid']
                    st.rerun()
        
        st.markdown(res)
        
        # Bot√µes de Download
        st.markdown("---")
        st.markdown("### üì• Download")
        col_down1, col_down2 = st.columns(2)
        
        with col_down1:
            docx_din = criar_docx_simples(res, "Din√¢mica Inclusiva")
            st.download_button(
                label="üìÑ Baixar DOCX",
                data=docx_din,
                file_name=f"Dinamica_{disciplina_bncc}_{date.today().strftime('%Y%m%d')}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        
        with col_down2:
            pdf_bytes_din = criar_pdf_generico(res)
            st.download_button(
                label="üìï Baixar PDF",
                data=pdf_bytes_din,
                file_name=f"Dinamica_{disciplina_bncc}_{date.today().strftime('%Y%m%d')}.pdf",
                mime="application/pdf",
                use_container_width=True
            )

def render_aba_plano_aula(aluno, api_key, kimi_key=None):
    """Renderiza a aba de plano de aula"""
    st.markdown("""
    <div class="pedagogia-box">
        <div class="pedagogia-title"><i class="ri-book-open-line"></i> Plano de Aula DUA</div>
        Gere um planejamento completo alinhado √† BNCC, selecionando metodologias ativas e recursos.
    </div>
    """, unsafe_allow_html=True)
    
    # BNCC em expander
    with st.expander("üìö BNCC e Habilidades", expanded=True):
        ano_bncc, disciplina_bncc, unidade_bncc, objeto_bncc, habilidades_bncc = criar_dropdowns_bncc_completos_melhorado(key_suffix="plano", aluno=aluno)
    
    # Configura√ß√£o Metodol√≥gica
    st.markdown("---")
    
    c1, c2 = st.columns(2)
    with c1:
        metodologia = st.selectbox("Metodologia", METODOLOGIAS, key="plano_met")
    
    tecnica_ativa = ""
    if metodologia == "Metodologia Ativa":
        with c2:
            tecnica_ativa = st.selectbox("T√©cnica Ativa", TECNICAS_ATIVAS, key="plano_tec")
    else:
        c2.info(f"Metodologia selecionada: {metodologia}")

    c3, c4 = st.columns(2)
    with c3:
        qtd_alunos_plano = st.number_input("Qtd Estudantes:", min_value=1, value=30, key="plano_qtd")
    with c4:
        recursos_plano = st.multiselect("Recursos Dispon√≠veis:", RECURSOS_DISPONIVEIS, key="plano_rec")
    
    # Bot√£o para gerar
    st.markdown("---")
    
    if st.button("üìÖ GERAR PLANO DE AULA", type="primary", use_container_width=True):
        if objeto_bncc and habilidades_bncc:
            with st.spinner(f"Consultando BNCC e planejando aula sobre '{objeto_bncc}'..."):
                res = gerar_plano_aula_completo(
                    api_key=api_key,
                    materia=disciplina_bncc,
                    assunto=objeto_bncc, # Passa o objeto BNCC como assunto
                    metodologia=metodologia,
                    tecnica=tecnica_ativa,
                    qtd_alunos=qtd_alunos_plano,
                    recursos=recursos_plano,
                    habilidades_bncc=habilidades_bncc,
                    verbos_bloom=None, # Bloom Removido
                    ano=ano_bncc,
                    unidade_tematica=unidade_bncc,
                    objeto_conhecimento=objeto_bncc,
                    aluno_info=aluno
                )
                st.session_state['res_plano'] = res
                st.session_state['res_plano_valid'] = False
        else:
            st.warning("‚ö†Ô∏è Selecione a 'Disciplina', 'Objeto do Conhecimento' e pelo menos uma 'Habilidade' nos campos da BNCC acima.")
    
    # Exibi√ß√£o do resultado
    if 'res_plano' in st.session_state:
        res = st.session_state['res_plano']
        
        st.markdown("---")
        st.markdown("### üìã Plano de Aula DUA")
        
        if st.session_state.get('res_plano_valid'):
            st.success("‚úÖ **PLANO VALIDADO E PRONTO PARA USO**")
        else:
            col_val, col_ajust, col_desc = st.columns(3)
            with col_val:
                if st.button("‚úÖ Validar Plano", key="val_plano", use_container_width=True):
                    st.session_state['res_plano_valid'] = True
                    ou.track_ai_feedback("hub", "validated", content_type="plano_aula")
                    st.rerun()
            with col_ajust:
                if st.button("üîÑ Refazer com Ajustes", key="redo_plano", use_container_width=True):
                    ou.track_ai_feedback("hub", "refazer", content_type="plano_aula")
                    st.session_state['res_plano_valid'] = False
                    st.info("Para ajustes, modifique os par√¢metros e clique em 'GERAR PLANO' novamente.")
            with col_desc:
                if st.button("üóëÔ∏è Descartar", key="del_plano", use_container_width=True):
                    del st.session_state['res_plano']
                    del st.session_state['res_plano_valid']
                    st.session_state.pop('res_plano_pptx', None)
                    st.rerun()
        
        st.markdown(res)
        
        # Bot√µes de Download
        st.markdown("---")
        st.markdown("### üì• Download")
        col_down1, col_down2, col_down3 = st.columns(3)
        
        with col_down1:
            docx_plano = criar_docx_simples(res, "Plano de Aula DUA")
            st.download_button(
                label="üìÑ Baixar DOCX",
                data=docx_plano,
                file_name=f"Plano_Aula_{disciplina_bncc}_{date.today().strftime('%Y%m%d')}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        
        with col_down2:
            pdf_bytes_plano = criar_pdf_generico(res)
            st.download_button(
                label="üìï Baixar PDF",
                data=pdf_bytes_plano,
                file_name=f"Plano_Aula_{disciplina_bncc}_{date.today().strftime('%Y%m%d')}.pdf",
                mime="application/pdf",
                use_container_width=True
            )
        
        with col_down3:
            ppt_key = "res_plano_pptx"
            if ppt_key not in st.session_state:
                st.session_state[ppt_key] = None
            if st.button(
                "üéØ Criar PPT (Omnisfera Green)",
                key="btn_ppt_plano",
                use_container_width=True,
                help="Gera PowerPoint a partir do plano usando Kimi (Omnisfera Green)"
            ):
                with st.spinner("Gerando PPT com Omnisfera Green..."):
                    ppt_bytes, err = gerar_ppt_do_plano_kimi(res, objeto_bncc or disciplina_bncc or "Plano de Aula", aluno=aluno, kimi_key=kimi_key)
                    if err:
                        st.session_state[ppt_key] = None
                        st.error(err)
                    else:
                        st.session_state[ppt_key] = ppt_bytes
                        st.rerun()
            if st.session_state.get(ppt_key):
                st.download_button(
                    label="üìä Baixar PPTX",
                    data=st.session_state[ppt_key],
                    file_name=f"Plano_Aula_{disciplina_bncc}_{date.today().strftime('%Y%m%d')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="dl_ppt_plano",
                    use_container_width=True
                )
    # ==============================================================================
# FUN√á√ïES DA EDUCA√á√ÉO INFANTIL
# ==============================================================================

def render_aba_ei_experiencia(aluno, api_key):
    """Renderiza a aba de experi√™ncia da Educa√ß√£o Infantil"""
    st.markdown("""
    <div class="pedagogia-box">
        <div class="pedagogia-title"><i class="ri-lightbulb-line"></i> Pedagogia do Brincar (BNCC)</div>
        Na Educa√ß√£o Infantil, n√£o fazemos "provas". Criamos <strong>experi√™ncias de aprendizagem</strong> intencionais. 
        Esta ferramenta usa a BNCC para criar brincadeiras que ensinam, usando o hiperfoco da crian√ßa.
    </div>
    """, unsafe_allow_html=True)
    
    col_ei1, col_ei2 = st.columns(2)
    campo_exp = col_ei1.selectbox("Campo de Experi√™ncia (BNCC)", CAMPOS_EXPERIENCIA_EI, key="campo_exp_ei")
    obj_aprendizagem = col_ei2.text_input("Objetivo de Aprendizagem:", placeholder="Ex: Compartilhar brinquedos, Identificar cores...", key="obj_aprendizagem_ei")
    
    if 'res_ei_exp' not in st.session_state:
        st.session_state.res_ei_exp = None
    if 'valid_ei_exp' not in st.session_state:
        st.session_state.valid_ei_exp = False

    if st.button("‚ú® GERAR EXPERI√äNCIA L√öDICA", type="primary", key="btn_exp_ei"):
        with st.spinner("Criando viv√™ncia..."):
            st.session_state.res_ei_exp = gerar_experiencia_ei_bncc(api_key, aluno, campo_exp=campo_exp, objetivo=obj_aprendizagem)
            st.session_state.valid_ei_exp = False

    if st.session_state.res_ei_exp:
        if st.session_state.valid_ei_exp:
            st.success("‚úÖ EXPERI√äNCIA APROVADA!")
            st.markdown(st.session_state.res_ei_exp)
        else:
            st.markdown(st.session_state.res_ei_exp)
            st.write("---")
            c_val, c_ref = st.columns([1, 3])
            if c_val.button("‚úÖ Validar Experi√™ncia", key="val_exp_ei"): 
                st.session_state.valid_ei_exp = True
                ou.track_ai_feedback("hub", "validated", content_type="experiencia_ei")
                st.rerun()
            with c_ref.expander("üîÑ N√£o gostou? Ensinar a IA"):
                feedback_ei = st.text_input("O que precisa melhorar?", placeholder="Ex: Ficou muito complexo, use materiais mais simples...", key="fb_exp_ei")
                if st.button("Refazer com Ajustes", key="refazer_exp_ei"):
                    ou.track_ai_feedback("hub", "refazer", content_type="experiencia_ei", feedback_text=feedback_ei or "")
                    with st.spinner("Reescrevendo..."):
                        st.session_state.res_ei_exp = gerar_experiencia_ei_bncc(api_key, aluno, campo_exp, obj_aprendizagem, feedback_anterior=feedback_ei)
                        st.rerun()

def render_aba_ei_estudio_visual(aluno, api_key, unsplash_key, gemini_key=None):
    """Renderiza a aba de est√∫dio visual da Educa√ß√£o Infantil (Gemini ou OpenAI + Unsplash)."""
    st.markdown(f"""
    <div class="pedagogia-box">
        <div class="pedagogia-title"><i class="ri-eye-line"></i> Apoio Visual & Comunica√ß√£o</div>
        Crian√ßas at√≠picas processam melhor imagens do que fala. 
        Use <strong>Cenas</strong> para hist√≥rias sociais e <strong>Pictogramas (CAA)</strong>. Usa <strong>{ou.AI_BLUE}</strong> ou <strong>{ou.AI_RED}</strong> (DALL-E)/Unsplash.
    </div>
    """, unsafe_allow_html=True)
    
    col_scene, col_caa = st.columns(2)
    g_key = gemini_key or ou.get_gemini_api_key()
    
    with col_scene:
        st.markdown("#### üñºÔ∏è Ilustra√ß√£o de Cena")
        desc_m = st.text_area("Descreva a cena ou rotina:", height=100, key="vdm_ei", placeholder="Ex: Crian√ßas brincando de roda no parque...")
        hiperfoco_ei = (aluno.get("hiperfoco") or "").strip()
        usar_hiperfoco_ei = st.checkbox(
            "Usar hiperfoco do estudante como tema da ilustra√ß√£o",
            value=bool(hiperfoco_ei),
            key="hub_usar_hiperfoco_ei",
            help="Se marcado, abre o campo com o hiperfoco para voc√™ editar antes de gerar.",
        )
        tema_ilustracao_ei = ""
        if usar_hiperfoco_ei:
            tema_ilustracao_ei = st.text_input(
                "Tema da ilustra√ß√£o (edite se quiser):",
                value=hiperfoco_ei,
                key="hub_tema_ilustracao_ei",
                placeholder="Ex: dinossauros, brincadeiras, animais...",
                help="Texto que ser√° usado como tema na gera√ß√£o. Vem do cadastro do estudante; voc√™ pode alterar.",
            )
            tema_ilustracao_ei = (tema_ilustracao_ei or "").strip()
        if st.button("üé® Gerar Cena", key="btn_cena_ei"):
            with st.spinner("Desenhando..."):
                prompt_completo = f"{desc_m}. Context: Child education, friendly style."
                if usar_hiperfoco_ei and tema_ilustracao_ei:
                    prompt_completo = f"Tema da ilustra√ß√£o: {tema_ilustracao_ei}. {prompt_completo}"
                st.session_state.res_scene_url = gerar_imagem_inteligente(api_key, prompt_completo, unsplash_key, prioridade="IA", gemini_key=g_key)
                st.session_state.valid_scene = False

        if st.session_state.res_scene_url:
            st.image(st.session_state.res_scene_url)
            if st.session_state.valid_scene:
                st.success("Imagem validada!")
            else:
                c_vs1, c_vs2 = st.columns([1, 2])
                if c_vs1.button("‚úÖ Validar", key="val_sc_ei"):
                    st.session_state.valid_scene = True
                    ou.track_ai_feedback("hub", "validated", content_type="cena_ei")
                    st.rerun()
                with c_vs2.expander("üîÑ Refazer Cena"):
                    fb_scene = st.text_input("Ajuste:", key="fb_sc_ei")
                    if st.button("Refazer", key="ref_sc_ei"):
                        ou.track_ai_feedback("hub", "refazer", content_type="cena_ei", feedback_text=fb_scene or "")
                        with st.spinner("Redesenhando..."):
                            prompt_completo = f"{desc_m}. Context: Child education."
                            if usar_hiperfoco_ei and tema_ilustracao_ei:
                                prompt_completo = f"Tema da ilustra√ß√£o: {tema_ilustracao_ei}. {prompt_completo}"
                            st.session_state.res_scene_url = gerar_imagem_inteligente(api_key, prompt_completo, unsplash_key, feedback_anterior=fb_scene, prioridade="IA", gemini_key=g_key)
                            st.rerun()

    with col_caa:
        st.markdown("#### üó£Ô∏è S√≠mbolo CAA (Comunica√ß√£o)")
        palavra_chave = st.text_input("Conceito/Palavra:", placeholder="Ex: Quero √Ågua, Banheiro", key="caa_input_ei")
        
        if st.button("üß© Gerar Pictograma", key="btn_caa_ei"):
            with st.spinner("Criando s√≠mbolo acess√≠vel..."):
                st.session_state.res_caa_url = gerar_pictograma_caa(api_key, palavra_chave, gemini_key=g_key)
                st.session_state.valid_caa = False

        if st.session_state.res_caa_url:
            st.image(st.session_state.res_caa_url, width=300)
            if st.session_state.valid_caa:
                st.success("Pictograma validado!")
            else:
                c_vc1, c_vc2 = st.columns([1, 2])
                if c_vc1.button("‚úÖ Validar", key="val_caa_ei"):
                    st.session_state.valid_caa = True
                    ou.track_ai_feedback("hub", "validated", content_type="pictograma_caa_ei")
                    st.rerun()
                with c_vc2.expander("üîÑ Refazer Picto"):
                    fb_caa = st.text_input("Ajuste:", key="fb_caa_ei")
                    if st.button("Refazer", key="ref_caa_ei"):
                        ou.track_ai_feedback("hub", "refazer", content_type="pictograma_caa_ei", feedback_text=fb_caa or "")
                        with st.spinner("Recriando..."):
                            st.session_state.res_caa_url = gerar_pictograma_caa(api_key, palavra_chave, feedback_anterior=fb_caa, gemini_key=g_key)
                            st.rerun()

def render_aba_ei_rotina(aluno, api_key):
    """Renderiza a aba de rotina da Educa√ß√£o Infantil"""
    st.markdown("""
    <div class="pedagogia-box">
        <div class="pedagogia-title"><i class="ri-calendar-check-line"></i> Rotina & Previsibilidade</div>
        A rotina organiza o pensamento da crian√ßa. Use esta ferramenta para identificar 
        pontos de estresse e criar estrat√©gias de antecipa√ß√£o.
    </div>
    """, unsafe_allow_html=True)
    
    rotina_detalhada = st.text_area("Descreva a Rotina da Turma:", height=200, placeholder="Ex: \n8:00 - Chegada e Acolhida\n8:30 - Roda de Conversa\n9:00 - Lanche\n...", key="rotina_ei")
    topico_foco = st.text_input("Ponto de Aten√ß√£o (Opcional):", placeholder="Ex: Transi√ß√£o para o parque", key="topico_foco_ei")
    
    if 'res_ei_rotina' not in st.session_state:
        st.session_state.res_ei_rotina = None
    if 'valid_ei_rotina' not in st.session_state:
        st.session_state.valid_ei_rotina = False

    if st.button("üìù ANALISAR E ADAPTAR ROTINA", type="primary", key="btn_rotina_ei"):
        with st.spinner("Analisando rotina..."):
            prompt_rotina = f"Analise esta rotina de Educa√ß√£o Infantil e sugira adapta√ß√µes sensoriais e visuais:\n\n{rotina_detalhada}\n\nFoco espec√≠fico: {topico_foco}"
            st.session_state.res_ei_rotina = gerar_roteiro_aula_completo(api_key, aluno, "Geral", "Rotina", feedback_anterior=prompt_rotina)
            st.session_state.valid_ei_rotina = False

    if st.session_state.res_ei_rotina:
        if st.session_state.valid_ei_rotina:
            st.success("‚úÖ ROTINA VALIDADA!")
            st.markdown(st.session_state.res_ei_rotina)
        else:
            st.markdown(st.session_state.res_ei_rotina)
            st.write("---")
            c_val, c_ref = st.columns([1, 3])
            if c_val.button("‚úÖ Validar Rotina", key="val_rotina_ei"):
                st.session_state.valid_ei_rotina = True
                ou.track_ai_feedback("hub", "validated", content_type="rotina_ei")
                st.rerun()
            with c_ref.expander("üîÑ Refazer adapta√ß√£o"):
                fb_rotina = st.text_input("O que ajustar na rotina?", key="fb_rotina_input_ei")
                if st.button("Refazer Rotina", key="refazer_rotina_ei"):
                    ou.track_ai_feedback("hub", "refazer", content_type="rotina_ei", feedback_text=fb_rotina or "")
                    with st.spinner("Reajustando..."):
                        prompt_rotina = f"Analise esta rotina de Educa√ß√£o Infantil e sugira adapta√ß√µes:\n\n{rotina_detalhada}\n\nFoco: {topico_foco}"
                        st.session_state.res_ei_rotina = gerar_roteiro_aula_completo(api_key, aluno, "Geral", "Rotina", feedback_anterior=prompt_rotina)
                        st.rerun()

def render_aba_ei_inclusao_brincar(aluno, api_key):
    """Renderiza a aba de inclus√£o no brincar da Educa√ß√£o Infantil"""
    st.markdown("""
    <div class="pedagogia-box">
        <div class="pedagogia-title"><i class="ri-group-line"></i> Media√ß√£o Social</div>
        Se a crian√ßa brinca isolada, o objetivo n√£o √© for√ßar a intera√ß√£o, mas criar 
        pontes atrav√©s do interesse dela. A IA criar√° uma brincadeira onde ela √© protagonista.
    </div>
    """, unsafe_allow_html=True)
    
    tema_d = st.text_input("Tema/Momento:", key="dina_ei", placeholder="Ex: Brincadeira de massinha")
    
    if 'res_ei_dina' not in st.session_state:
        st.session_state.res_ei_dina = None
    if 'valid_ei_dina' not in st.session_state:
        st.session_state.valid_ei_dina = False

    if st.button("ü§ù GERAR DIN√ÇMICA", type="primary", key="btn_dina_ei"): 
        with st.spinner("Criando ponte social..."):
            st.session_state.res_ei_dina = gerar_dinamica_inclusiva_completa(
                api_key, 
                aluno, 
                "Educa√ß√£o Infantil", 
                tema_d, 
                10,  # pequeno grupo
                "Crian√ßas pequenas"
            )
            st.session_state.valid_ei_dina = False

    if st.session_state.res_ei_dina:
        if st.session_state.valid_ei_dina:
            st.success("‚úÖ DIN√ÇMICA VALIDADA!")
            st.markdown(st.session_state.res_ei_dina)
        else:
            st.markdown(st.session_state.res_ei_dina)
            st.write("---")
            c_val, c_ref = st.columns([1, 3])
            if c_val.button("‚úÖ Validar Din√¢mica", key="val_dina_ei"):
                st.session_state.valid_ei_dina = True
                ou.track_ai_feedback("hub", "validated", content_type="dinamica_ei")
                st.rerun()
            with c_ref.expander("üîÑ Refazer din√¢mica"):
                fb_dina = st.text_input("O que ajustar?", key="fb_dina_input_ei")
                if st.button("Refazer Din√¢mica", key="refazer_dina_ei"):
                    ou.track_ai_feedback("hub", "refazer", content_type="dinamica_ei", feedback_text=fb_dina or "")
                    with st.spinner("Reajustando..."):
                        st.session_state.res_ei_dina = gerar_dinamica_inclusiva_completa(
                            api_key, 
                            aluno, 
                            "Educa√ß√£o Infantil", 
                            tema_d, 
                            10, 
                            "Crian√ßas pequenas", 
                            feedback_anterior=fb_dina
                        )
                        st.rerun()

# ==============================================================================
# EXECU√á√ÉO PRINCIPAL
# ==============================================================================

# CSS espec√≠fico do Hub (chamado ap√≥s a defini√ß√£o da fun√ß√£o)
aplicar_estilos()

# ==============================================================================
# FUN√á√ÉO PRINCIPAL
# ==============================================================================
def main():
    """Fun√ß√£o principal da aplica√ß√£o - executa a l√≥gica do Hub"""
    
    # Inicializar api_key e unsplash_key antes de usar (Render: env var ‚Üí Streamlit Cloud: secrets ‚Üí sess√£o)
    api_key = os.environ.get("OPENAI_API_KEY") or ou.get_setting("OPENAI_API_KEY", "") or st.session_state.get("OPENAI_API_KEY", "")
    api_key = (api_key or "").strip() or None
    
    unsplash_key = (
        os.environ.get("UNSPLASH_ACCESS_KEY")
        or ou.get_setting("UNSPLASH_ACCESS_KEY", "")
        or st.session_state.get("UNSPLASH_ACCESS_KEY")
    )
    unsplash_key = unsplash_key if unsplash_key else None
    
    gemini_key = ou.get_gemini_api_key()
    kimi_key = (
        st.session_state.get("_kimi_key_cached")
        or os.environ.get("OPENROUTER_API_KEY")
        or os.environ.get("KIMI_API_KEY")
        or ou.get_setting("OPENROUTER_API_KEY", "")
        or ou.get_setting("KIMI_API_KEY", "")
        or st.session_state.get("OPENROUTER_API_KEY", "")
        or st.session_state.get("KIMI_API_KEY", "")
    )
    try:
        kimi_key = kimi_key or (st.secrets.get("OPENROUTER_API_KEY", "") or st.secrets.get("KIMI_API_KEY", "") or "")
    except Exception:
        pass
    kimi_key = (kimi_key or "").strip() or None
    if kimi_key:
        st.session_state["_kimi_key_cached"] = kimi_key

    # Carregar dados dos estudantes
    if 'banco_estudantes' not in st.session_state or not st.session_state.banco_estudantes:
        with st.spinner("üîÑ Lendo dados da nuvem..."):
            st.session_state.banco_estudantes = carregar_estudantes_supabase()
    
    if not st.session_state.banco_estudantes:
        st.warning("‚ö†Ô∏è Nenhum estudante encontrado.")
        try:
            from ui.permissions import get_member_from_session
            from services.members_service import get_class_assignments
            m = get_member_from_session()
            if m and (m.get("link_type") or "").lower() == "turma":
                ca = get_class_assignments(m.get("id", ""))
                if not ca:
                    st.info("üí° Voc√™ est√° vinculado **por turma** mas n√£o tem turmas atribu√≠das. Pe√ßa √† coordena√ß√£o para configurar em **Gest√£o de Usu√°rios** (turmas e componentes).")
                else:
                    st.info("üí° Nenhum aluno corresponde √†s suas turmas no momento. Verifique se o PEI dos alunos est√° com a mesma s√©rie/turma das suas atribui√ß√µes.")
        except Exception:
            pass
        if st.button("üìò Ir para o m√≥dulo PEI", type="primary"):
            st.switch_page("pages/1_PEI.py")
        st.stop()
    
    # Sele√ß√£o de aluno + Bot√£o de Limpar Formul√°rios
    lista_alunos = [a['nome'] for a in st.session_state.banco_estudantes]
    col_sel, col_btn_limpar, col_info = st.columns([2, 1, 2])
    with col_sel:
        nome_aluno = st.selectbox("üìÇ Selecione o Estudante:", lista_alunos)
    
    with col_btn_limpar:
        st.markdown("<div style='height:28px'></div>", unsafe_allow_html=True)  # Alinhar com selectbox
        if st.button("üßπ Limpar Formul√°rios", type="secondary", use_container_width=True, help="Limpa todos os formul√°rios e resultados da p√°gina atual"):
            # Limpar todos os estados relacionados a formul√°rios e resultados
            chaves_para_limpar = [
                # Resultados de adapta√ß√£o
                'res_docx', 'res_img', 'res_create', 
                # Valida√ß√µes
                'valid_scene', 'valid_caa', 'valid_d', 'valid_i', 'valid_c',
                # Uploads e arquivos
                'docx_txt', 'docx_imgs', 'last_d', 'last_i', 'img_raw', 'cropped_img',
                # Educa√ß√£o Infantil
                'res_ei_experiencia', 'valid_ei_experiencia', 'res_ei_rotina', 
                'valid_ei_rotina', 'res_ei_dina', 'valid_ei_dina',
                # Outras abas
                'res_roteiro', 'res_roteiro_valid', 'res_papo', 'res_papo_valid',
                'res_dinamica', 'res_dinamica_valid', 'res_plano_aula', 'res_plano_aula_valid',
                'res_plano', 'res_plano_valid', 'res_plano_pptx',
                # Configura√ß√µes de adapta√ß√£o
                'necessidades_selecionadas_adaptar_prova',
                # Estados de Bloom (limpar todos os prefixos)
            ]
            
            # Limpar chaves espec√≠ficas
            for key in chaves_para_limpar:
                if key in st.session_state:
                    del st.session_state[key]
            
            # Limpar estados de Bloom (todos os prefixos)
            keys_to_remove = []
            for key in st.session_state.keys():
                if key.startswith('bloom_memoria_') or key.startswith('usar_bloom_') or key.startswith('cat_bloom_') or key.startswith('ms_bloom_'):
                    keys_to_remove.append(key)
            for key in keys_to_remove:
                del st.session_state[key]
            
            st.success("‚úÖ Formul√°rios limpos!")
            st.rerun()
    
    aluno = next((a for a in st.session_state.banco_estudantes if a.get('nome') == nome_aluno), None)
    
    if not aluno: 
        st.error("Estudante n√£o encontrado")
        st.stop()
    
    # --- √ÅREA DO ALUNO (Visual + Expander) ---
    
    # 1. Cabe√ßalho Visual (Nome e S√©rie - Mant√©m vis√≠vel para contexto r√°pido)
    render_cabecalho_aluno(aluno)
    
    # 2. PEI RETR√ÅTIL (Aqui est√° a mudan√ßa solicitada)
    # Mostra os detalhes pesados apenas se o usu√°rio clicar
    with st.expander(f"üìÑ Ver Detalhes do PEI ‚Äî {aluno['nome'].split()[0]} (uso interno da equipe)", expanded=False):
        st.caption("Contexto pedag√≥gico e cl√≠nico. Estas informa√ß√µes permanecem na Omnisfera e n√£o s√£o compartilhadas com o estudante ou a fam√≠lia.")
        # Buscar diagn√≥stico - IMPORTANTE: usar diretamente do aluno que j√° foi processado
        # O diagn√≥stico j√° foi extra√≠do corretamente na fun√ß√£o carregar_estudantes_supabase
        diagnostico_pei = aluno.get('diagnosis', '') or ""
        
        # Se n√£o encontrou ou est√° vazio, tenta buscar diretamente do pei_data como fallback
        if not diagnostico_pei or diagnostico_pei == "N√£o informado no cadastro":
            pei_data_local = aluno.get('pei_data', {}) or {}
            if isinstance(pei_data_local, dict):
                diagnostico_pei = pei_data_local.get('diagnostico', '') or ""
        
        # Se ainda n√£o encontrou, mostra mensagem
        if not diagnostico_pei:
            diagnostico_pei = "N√£o informado no cadastro"
        
        # Buscar hiperfoco - IMPORTANTE: usar diretamente do aluno que j√° foi processado
        # O hiperfoco j√° foi extra√≠do corretamente do pei_data na fun√ß√£o carregar_estudantes_supabase
        hiperfoco_aluno = aluno.get('hiperfoco', '') or ""
        
        # Se n√£o encontrou, tenta buscar diretamente do pei_data como fallback
        if not hiperfoco_aluno or hiperfoco_aluno == "Interesses gerais (A descobrir)":
            pei_data_local = aluno.get('pei_data', {}) or {}
            if isinstance(pei_data_local, dict):
                hiperfoco_temp = (
                    pei_data_local.get('hiperfoco') or 
                    pei_data_local.get('interesses') or 
                    pei_data_local.get('habilidades_interesses') or 
                    ""
                )
                # Garantir que o hiperfoco n√£o seja igual ao diagn√≥stico
                if hiperfoco_temp and hiperfoco_temp != diagnostico_pei:
                    hiperfoco_aluno = hiperfoco_temp
        
        # Se ainda n√£o encontrou hiperfoco, mostra padr√£o
        if not hiperfoco_aluno or hiperfoco_aluno == diagnostico_pei:
            hiperfoco_aluno = "Interesses gerais (A descobrir)"
        
        c_diag, c_hip = st.columns(2)
        with c_diag:
            st.markdown("**üè• Contexto cl√≠nico (apenas equipe):**")
            # Garantir que est√° usando o campo correto
            diag_final = diagnostico_pei if diagnostico_pei and diagnostico_pei != "N√£o informado no cadastro" else "N√£o informado no cadastro"
            st.info(diag_final)
            
        with c_hip:
            st.markdown("**üéØ Interesses / Hiperfoco:**")
            # Garantir que est√° usando o campo correto - N√ÉO usar diagn√≥stico aqui
            hip_final = hiperfoco_aluno if hiperfoco_aluno and hiperfoco_aluno != "N√£o informado no cadastro" else "Interesses gerais (A descobrir)"
            st.success(hip_final)
            
        st.markdown("---")
        st.markdown("**üß† Resumo das Estrat√©gias (IA):**")
        # Mostra o texto da IA (ia_sugestao) formatado
        st.write(aluno.get('ia_sugestao', 'Sem resumo dispon√≠vel.'))

    # --- FIM DA √ÅREA DO ALUNO ---

    # Detector de Educa√ß√£o Infantil
    serie_aluno = aluno.get('serie', '').lower()
    is_ei = any(termo in serie_aluno for termo in ["infantil", "creche", "pr√©", "pre", "maternal", "ber√ß√°rio"])
    
    # Inicializar estado para imagens
    if 'res_scene_url' not in st.session_state:
        st.session_state.res_scene_url = None
    if 'valid_scene' not in st.session_state:
        st.session_state.valid_scene = False
    if 'res_caa_url' not in st.session_state:
        st.session_state.res_caa_url = None
    if 'valid_caa' not in st.session_state:
        st.session_state.valid_caa = False
    
    if is_ei:
        # Modo Educa√ß√£o Infantil
        st.info("üß∏ **Modo Educa√ß√£o Infantil Ativado:** Foco em Experi√™ncias, BNCC e Brincar.")
        
        tabs = st.tabs([
            "üß∏ Criar Experi√™ncia (BNCC)",
            "üé® Est√∫dio Visual & CAA",
            "üìù Rotina & AVD",
            "ü§ù Inclus√£o no Brincar",
        ])
        
        with tabs[0]:
            render_aba_ei_experiencia(aluno, api_key)
        
        with tabs[1]:
            render_aba_ei_estudio_visual(aluno, api_key, unsplash_key, gemini_key)
        
        with tabs[2]:
            render_aba_ei_rotina(aluno, api_key)
        
        with tabs[3]:
            render_aba_ei_inclusao_brincar(aluno, api_key)
    
    else:
        # Modo Padr√£o (Fundamental / M√©dio)
        tabs = st.tabs([
            "üìÑ Adaptar Prova",
            "‚úÇÔ∏è Adaptar Atividade",
            "‚ú® Criar do Zero",
            "üé® Est√∫dio Visual & CAA",
            "üìù Roteiro Individual",
            "üó£Ô∏è Papo de Mestre",
            "ü§ù Din√¢mica Inclusiva",
            "üìÖ Plano de Aula DUA",
        ])
        
        with tabs[0]:
            render_aba_adaptar_prova(aluno, api_key)
        
        with tabs[1]:
            render_aba_adaptar_atividade(aluno, api_key)
        
        with tabs[2]:
            render_aba_criar_do_zero(aluno, api_key, unsplash_key)

        with tabs[3]:
            render_aba_estudio_visual(aluno, api_key, unsplash_key, gemini_key)

        with tabs[4]:
            render_aba_roteiro_individual(aluno, api_key)

        with tabs[5]:
            render_aba_papo_mestre(aluno, api_key)

        with tabs[6]:
            render_aba_dinamica_inclusiva(aluno, api_key)

        with tabs[7]:
            render_aba_plano_aula(aluno, api_key, kimi_key)

# Executa a fun√ß√£o principal
main()

# ==============================================================================
# RODAP√â COM ASSINATURA
# ==============================================================================
ou.render_footer_assinatura()








    
